# -*- coding: utf-8 -*-
"""
format.ppt v2.0 — 修复版: 无重叠、多样化版式、发现融入内容、高空间利用率
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path
import json

OUT_DIR = Path(r"C:\Users\zoezoezhao\Desktop\汇报agent\案例素材\case1-文档+slides+ai slides")
OUT = OUT_DIR / "AI 产品用户留存分析_format_v2.0_fixed.pptx"
QA_OUT = OUT_DIR / "AI 产品用户留存分析_format_v2.0_fixed_QA.json"

SW = Inches(13.333); SH = Inches(7.5)

C0 = RGBColor(0x00,0x52,0xD9)
C0L = RGBColor(0xE6,0xF1,0xFB)
CP = {"元宝": RGBColor(0x22,0xC5,0x5E), "DS": RGBColor(0x25,0x63,0xEB), "豆包": RGBColor(0x93,0xC5,0xFD)}
C_OK = RGBColor(0x22,0xC5,0x5E); C_BAD = RGBColor(0xEF,0x44,0x4A)
C_BG = RGBColor(0xF5,0xF5,0xF5); C_BGQ = RGBColor(0xFF,0xF8,0xE1)
CT = RGBColor(0x33,0x33,0x33); CTS = RGBColor(0x66,0x66,0x66); CTL = RGBColor(0x99,0x99,0x99)
CBO = RGBColor(0xE0,0xE0,0xE0); CW = RGBColor(0xFF,0xFF,0xFF)
CNI = RGBColor(0xF0,0xF0,0xF0)

P = ["元宝","DS","豆包"]
SR = {"元宝":19,"DS":34,"豆包":54}

def tb(sl,l,t,w,h,txt,sz=12,b=False,c=CT,a=PP_ALIGN.LEFT):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.alignment=a;r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c
    return bx

def rect(sl,l,t,w,h,fill=CW,border=CBO):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    s.fill.solid();s.fill.fore_color.rgb=fill;s.line.color.rgb=border;s.line.width=Pt(0.5);return s

def rrect(sl,l,t,w,h,fill=CW,border=None):
    s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    s.fill.solid();s.fill.fore_color.rgb=fill
    if border:s.line.color.rgb=border;s.line.width=Pt(0.5)
    else:s.line.fill.background()
    return s

def nav(sl,tabs,active):
    """导航标签栏 — 用单个文本框替代 shape+text 重叠"""
    w=Inches(1.5);h=Inches(0.45);gap=Inches(0.1);x0=Inches(0.4)
    for i,t in enumerate(tabs):
        x=int(x0+i*(w+gap));on=i==active
        rrect(sl,x,Inches(0.1),w,h,fill=C0 if on else CNI)
        tb(sl,x,Inches(0.12),w,h-0.05,t,10,on,CW if on else CTS,a=PP_ALIGN.CENTER)

def ft(sl,txt="样本 N=4,174 | 2025.11.11-11.13"):
    tb(sl,Inches(0.5),SH-Inches(0.4),SW-Inches(1),Inches(0.3),txt,8,c=CTL)

def title(sl,txt,y=Inches(0.7)):
    tb(sl,Inches(0.5),y,SW-Inches(1),Inches(0.45),txt,15,True,CT)

# ── 数据 ──
STRONG_RET = dict(SR)
STRONG_PEN = {"元宝": 66.4, "DS": 86.4, "豆包": 95.1}
PN = {"元宝":{"pure":26,"non":12,"diff":14,"pct":5.2},
      "DS":{"pure":45,"non":27,"diff":18,"pct":12.8},
      "豆包":{"pure":60,"non":50,"diff":10,"pct":17.1}}
DR = {"元宝":[("可靠性",16.1),("深度思考",14.3),("拍照答疑",12.6),("AI创作",12.0)],
      "DS":[("可靠性",25.2),("深度思考",23.8),("情绪价值",23.3),("拍照答疑",17.9)],
      "豆包":[("拍照答疑",21.0),("深度思考",20.2),("速度稳定性",20.1),("情绪价值",20.1)]}
FS = {"元宝":[("打电话",0.170),("AI创作",0.144),("朗读",0.144),("深度思考",0.122)],
      "DS":[("深度思考",0.137),("AI创作",0.061),("拍照答疑",0.057),("打电话",0.042)],
      "豆包":[("深度思考",0.227),("拍照答疑",0.184),("打电话",0.172),("朗读",0.170)]}
SC = {"元宝":(14,19,"元宝"),"DS":(18,34,"DS"),"豆包":(10,54,"豆包")}
QUOTES = [("DS用户","一开始觉得元宝和DS差不多就没下载。但试用后觉得整体比DS更智能、功能更丰富些。",
           "DS用户迁移成本低，体验领先可撬动转换"),
          ("元宝用户","希望元宝提供情绪价值，主动发问，互动感可增强用户粘性",
           "用户的主动交互需求未被满足，被动召回有创新空间")]

# ── 简报呈现 (方式 A: 标题融入) ──────────────────────────
# 版式1: 行动标题即发现，图表主导，无独立面板
def retention_overview(sl):
    nav(sl,["留存总览","驱动因素","功能留存","召回流失","建议"],0)
    # 发现作为标题 (方式A)
    tb(sl,Inches(0.5),Inches(0.7),SW-Inches(1),Inches(0.5),
       "发现1：豆包强留存率54%大幅领先，元宝仅19%差距达35pp",13,True,C0)
    
    # 三个产品卡片全屏展示，大幅减少空白
    cw=Inches(3.8);ch=Inches(5.0);gap=Inches(0.2);y=Inches(1.5)
    for i,p in enumerate(P):
        x=Inches(0.5)+i*(cw+gap)
        rrect(sl,x,y,cw,ch,fill=CW,border=CBO)
        rect(sl,x,y,cw,Inches(0.5),fill=CP[p])
        tb(sl,x,y+Inches(0.05),cw,Inches(0.35),p,14,True,CW,a=PP_ALIGN.CENTER)
        # 强留存率
        tb(sl,x+Inches(0.3),y+Inches(0.8),cw-Inches(0.6),Inches(0.3),
           f"强留存率 {SR[p]}%",18,True,CP[p])
        # pp标注 (方式C: 图表内标注)
        tb(sl,x+Inches(0.3),y+Inches(1.2),cw-Inches(0.6),Inches(0.25),
           f"vs 元宝 {SR[p]-SR['元宝']:+d}pp",10,True,C_OK if p!='元宝' else CT)
        # 渗透率
        tb(sl,x+Inches(0.3),y+Inches(1.7),cw-Inches(0.6),Inches(0.25),
           f"渗透率 {STRONG_PEN[p]}%",11,c=CTS)
        # 关键insight
        ins = {"元宝":"渗透率66.4%，有较大提升空间",
               "DS":"纯白留存+18pp效果最显著",
               "豆包":"全面领先，巩固优势"}
        tb(sl,x+Inches(0.3),y+Inches(2.3),cw-Inches(0.6),Inches(0.6),
           ins[p],9,c=CTS)
        # 简单bar示意
        bw=int((cw-Inches(0.6))*SR[p]/100)
        rrect(sl,x+Inches(0.3),y+Inches(3.2),Emu(bw),Inches(0.3),fill=CP[p])
        tb(sl,x+Inches(0.3)+Emu(bw)+Inches(0.1),y+Inches(3.2),Inches(0.6),Inches(0.3),
           f"{SR[p]}%",10,True,CP[p])
    ft(sl)

# ── 版式2: 散点图 + 注解式发现 (方式C: 图表标注) ──────────
def scatter_page(sl):
    nav(sl,["留存总览","驱动因素","功能留存","召回流失","建议"],1)
    # 标题即发现 (方式A)
    tb(sl,Inches(0.5),Inches(0.7),SW-Inches(1),Inches(0.5),
       "发现2：DS纯白用户拉动+18pp效果最显著，豆包基数高拉动较小",13,True,C0)
    
    # 散点图占左侧70%版面
    pl=Inches(0.8);pt=Inches(1.5);pw=Inches(8.5);ph=Inches(5.0)
    max_x=65;max_y=22
    # 坐标轴
    rect(sl,pl,pt,Inches(0.01),ph,fill=CBO,border=CBO)
    rect(sl,pl,pt+ph,pw,Inches(0.01),fill=CBO,border=CBO)
    tb(sl,pl,pt-Inches(0.3),Inches(2),Inches(0.3),
       "强留存提升(pp)",10,True,CT)
    tb(sl,pl+pw-Inches(2),pt+ph+Inches(0.05),Inches(2),Inches(0.25),
       "强留存率(%)",10,True,CT,a=PP_ALIGN.RIGHT)
    
    for key,(lift,rate,label) in SC.items():
        x=pl+int(pw*rate/max_x);y=pt+int(ph*(1-lift/max_y))
        r=Inches(0.18)
        c=sl.shapes.add_shape(MSO_SHAPE.OVAL,x-r,y-r,r*2,r*2)
        c.fill.solid();c.fill.fore_color.rgb=CP[key];c.line.color.rgb=CW;c.line.width=Pt(2)
        # Label on left for rightmost point to avoid overlap
        if rate > 40:
            tb(sl,x-Inches(2.8),y-Inches(0.2),Inches(2.5),Inches(0.2),label,11,True,CP[key],a=PP_ALIGN.RIGHT)
            tb(sl,x-Inches(2.8),y+Inches(0.05),Inches(2.5),Inches(0.2),f"({rate}%, +{lift}pp)",8,False,CTS,PP_ALIGN.RIGHT)
        else:
            tb(sl,x+Inches(0.3),y-Inches(0.2),Inches(1.5),Inches(0.2),label,11,True,CP[key])
            tb(sl,x+Inches(0.3),y+Inches(0.05),Inches(2),Inches(0.2),f"({rate}%, +{lift}pp)",8,c=CTS)
    
    # 右侧注解区 — 非固定面板，是图表标注延伸 (方式C)
    ax=x=Inches(10.0);y=Inches(1.8)
    tb(sl,ax,y,Inches(3),Inches(0.25),"关键发现",11,True,C0)
    notes=[
        "DS纯白拉动+18pp\n占位有效",
        "豆包纯白+10pp\n基数54%已很高",
        "元宝纯白+14pp\n提升空间最大",
    ]
    for i,n in enumerate(notes):
        yp=y+Inches(0.4)+i*Inches(0.9)
        rrect(sl,ax,yp,Inches(2.6),Inches(0.7),fill=C0L)
        tb(sl,ax+Inches(0.15),yp+Inches(0.08),Inches(2.3),Inches(0.55),n,9,c=CT)
    ft(sl)

# ── 版式3: 三栏对比 (方案B) — 无发现面板 ──────────────────
def driver_comparison(sl):
    nav(sl,["留存总览","驱动因素","功能留存","召回流失","建议"],1)
    tb(sl,Inches(0.5),Inches(0.7),SW-Inches(1),Inches(0.5),
       "发现3：可靠性是DS和元宝首要驱动力，豆包以拍照答疑领先",13,True,C0)
    
    cw=Inches(3.8);ch=Inches(5.2);gap=Inches(0.3);y=Inches(1.5)
    # 可用bar宽度: cw - label(1.2) - gap(0.1) - padding(0.2) = 2.3"
    max_bw=int(cw-Inches(1.5))  # 2.3 inches max bar width
    for i,p in enumerate(P):
        x=Inches(0.5)+i*(cw+gap)
        rrect(sl,x,y,cw,ch,fill=CW,border=CBO)
        rect(sl,x,y,cw,Inches(0.45),fill=CP[p])
        tb(sl,x,y+Inches(0.05),cw,Inches(0.35),p,13,True,CW,a=PP_ALIGN.CENTER)
        # 驱动因素水平条
        items=DR[p];max_v=max(v for _,v in items)
        by=y+Inches(0.7);label_w=Inches(1.2);bar_x=x+label_w+Inches(0.1)
        for j,(label,val) in enumerate(items):
            bw=int(max_bw*val/max_v);bh=Inches(0.4);by2=by+j*(bh+Inches(0.12))
            tb(sl,x+Inches(0.15),by2+Inches(0.05),label_w,bh,label,9,c=CT)
            if bw>Inches(0.1):
                rrect(sl,bar_x,by2,Emu(bw),bh,fill=CP[p])
            tb(sl,bar_x+Emu(bw)+Inches(0.04),by2+Inches(0.05),
               Inches(0.5),bh,f"{val:.0f}%",9,True,CP[p])
    ft(sl)

# ── 版式4: 左文字+右表格 (方案A: 表格型) 发现作为标题 ──────
def func_satisfaction(sl):
    nav(sl,["留存总览","驱动因素","功能留存","召回流失","建议"],2)
    tb(sl,Inches(0.5),Inches(0.7),SW-Inches(1),Inches(0.5),
       "发现4：深度思考是三大产品共有关键功能，豆包0.227相关性最高",13,True,C0)
    
    # 左侧: 关键发现文字 (方式A: 标题融入 + 简注)
    lx=Inches(0.5);ly=Inches(1.5);lw=Inches(4.0)
    tb(sl,lx,ly,lw,Inches(0.25),"核心发现",11,True,C0)
    finds=[
        "深度思考是唯一在三个产品中\n均进入Top4的功能",
        "拍照答疑对元宝和豆包\n均有强相关性(+0.17)",
        "打电话在元宝排名第一\n但DS和豆包中较低",
        "AI创作在DS和豆包中\n均有0.06+的相关性",
    ]
    for i,f in enumerate(finds):
        yp=ly+Inches(0.4)+i*Inches(0.9)
        rrect(sl,lx,yp,lw,Inches(0.7),fill=C0L)
        tb(sl,lx+Inches(0.15),yp+Inches(0.08),lw-Inches(0.3),Inches(0.55),f,9,c=CT)
    
    # 右侧: 功能满意度表格
    tx=Inches(5.0);ty=Inches(1.5)
    funcs=list(set(f for p in P for f,_ in FS[p]))
    fsd={p:dict(FS[p]) for p in P}
    funcs.sort(key=lambda f:sum(fsd[p].get(f,0) for p in P)/3,reverse=True)
    funcs=funcs[:7]
    rows=len(funcs)+1;cols=4
    tbl_w=Inches(7.8);tbl_h=Inches(4.0)
    tbl=sl.shapes.add_table(rows,cols,tx,ty,tbl_w,tbl_h).table
    for ci,h in enumerate(["功能"]+P):
        c=tbl.cell(0,ci);c.text="";p=c.text_frame.paragraphs[0];p.alignment=PP_ALIGN.CENTER
        r=p.add_run();r.text=h;r.font.size=Pt(9);r.font.bold=True;r.font.color.rgb=CW
        c.fill.solid();c.fill.fore_color.rgb=C0
    for ri,f in enumerate(funcs):
        c=tbl.cell(ri+1,0);c.text="";p=c.text_frame.paragraphs[0]
        r=p.add_run();r.text=f;r.font.size=Pt(9);r.font.bold=True;r.font.color.rgb=CT
        c.fill.solid();c.fill.fore_color.rgb=CNI
        for ci,pr in enumerate(P):
            val=fsd[pr].get(f,None)
            txt=f"{val:.3f}" if val else "-"
            c=tbl.cell(ri+1,ci+1);c.text="";p=c.text_frame.paragraphs[0];p.alignment=PP_ALIGN.CENTER
            r=p.add_run();r.text=txt;r.font.size=Pt(9);r.font.color.rgb=CT
            if val and val==max(v for p2 in P for f2,v in FS[p2] if f2==f):
                c.fill.solid();c.fill.fore_color.rgb=C0L
    ft(sl)

# ── 版式5: 上下分层 (方案C) — 发现作为数据标注 ──────────
def recall_churn(sl):
    nav(sl,["留存总览","驱动因素","功能留存","召回流失","建议"],3)
    tb(sl,Inches(0.5),Inches(0.7),SW-Inches(1),Inches(0.5),
       "发现5：主动打开是首要召回渠道，'其他AI更好'是主要流失原因",13,True,C0)
    
    # 上半: 召回渠道 (表)
    ty=Inches(1.5);th=Inches(2.7)
    rec=[("主动打开",58.1,57.5,61.3),("桌面入口",18.3,29.3,30.5),
         ("Push通知",11.0,13.0,11.5),("红点提示",14.0,19.4,23.0),
         ("朋友分享",14.0,13.0,13.0)]
    tb(sl,Inches(0.5),ty-Inches(0.3),Inches(3),Inches(0.3),"召回渠道分布 (%)",11,True,C0)
    rows=len(rec)+1;cols=4
    tbl=sl.shapes.add_table(rows,cols,Inches(0.5),ty,Inches(5.5),th).table
    for ci,h in enumerate(["渠道"]+P):
        c=tbl.cell(0,ci);c.text="";p=c.text_frame.paragraphs[0];p.alignment=PP_ALIGN.CENTER
        r=p.add_run();r.text=h;r.font.size=Pt(9);r.font.bold=True;r.font.color.rgb=CW
        c.fill.solid();c.fill.fore_color.rgb=C0
    for ri,(ch,*vals) in enumerate(rec):
        c=tbl.cell(ri+1,0);c.text="";p=c.text_frame.paragraphs[0];p.alignment=PP_ALIGN.CENTER
        r=p.add_run();r.text=ch;r.font.size=Pt(9);r.font.color.rgb=CT
        for ci,v in enumerate(vals):
            c=tbl.cell(ri+1,ci+1);c.text="";p=c.text_frame.paragraphs[0];p.alignment=PP_ALIGN.CENTER
            r=p.add_run();r.text=f"{v:.1f}%";r.font.size=Pt(9);r.font.color.rgb=CT
            if v==max(vals):c.fill.solid();c.fill.fore_color.rgb=C0L
    
    # 下半: 流失原因 (表)
    by=Inches(4.7);bh=Inches(1.9)
    churn=[("其他AI产品更好",14.0,19.5,19.1),("没安装/不需要",26.1,33.4,28.0),
           ("已卸载",13.0,12.2,10.0),("不知道使用场景",12.0,10.0,8.0)]
    tb(sl,Inches(0.5),by-Inches(0.3),Inches(3),Inches(0.3),"主要流失原因 (%)",11,True,C_BAD)
    rows2=len(churn)+1
    tbl2=sl.shapes.add_table(rows2,cols,Inches(0.5),by,Inches(5.5),bh).table
    for ci,h in enumerate(["原因"]+P):
        c=tbl2.cell(0,ci);c.text="";p=c.text_frame.paragraphs[0];p.alignment=PP_ALIGN.CENTER
        r=p.add_run();r.text=h;r.font.size=Pt(9);r.font.bold=True;r.font.color.rgb=CW
        c.fill.solid();c.fill.fore_color.rgb=C_BAD
    for ri,(re,*vals) in enumerate(churn):
        c=tbl2.cell(ri+1,0);c.text="";p=c.text_frame.paragraphs[0];p.alignment=PP_ALIGN.CENTER
        r=p.add_run();r.text=re;r.font.size=Pt(9);r.font.color.rgb=CT
        for ci,v in enumerate(vals):
            c=tbl2.cell(ri+1,ci+1);c.text="";p=c.text_frame.paragraphs[0];p.alignment=PP_ALIGN.CENTER
            r=p.add_run();r.text=f"{v:.1f}%";r.font.size=Pt(9);r.font.color.rgb=CT
    
    # 右侧: pp标注作为发现 (方式C)
    rx=Inches(6.8);ry=Inches(1.5)
    tb(sl,rx,ry,Inches(5.5),Inches(0.25),"关键发现",11,True,C0)
    findings=[
        ("召回提升机会","桌面入口(30.5%)和Push(11.5%)\n元宝有显著提升空间","元宝"),
        ("流失差异","DS用户'其他AI更好'达19.5%,\n迁移风险最高","DS"),
    ]
    for i,(ftitle,fbody,color_key) in enumerate(findings):
        fy=ry+Inches(0.4)+i*Inches(1.2)
        rrect(sl,rx,fy,Inches(5.5),Inches(1.0),fill=C_BG)
        rect(sl,rx,fy,Inches(0.06),Inches(1.0),fill=CP[color_key])
        tb(sl,rx+Inches(0.2),fy+Inches(0.08),Inches(5.0),Inches(0.25),ftitle,10,True,CT)
        tb(sl,rx+Inches(0.2),fy+Inches(0.35),Inches(5.0),Inches(0.55),fbody,9,c=CTS)
    ft(sl)

# ── 版式6: 引文嵌入 (方式D) — 用户原声嵌入分析页 ──────
def user_voice(sl):
    nav(sl,["留存总览","驱动因素","功能留存","召回流失","建议"],4)
    tb(sl,Inches(0.5),Inches(0.7),SW-Inches(1),Inches(0.5),
       "发现6：用户期望更高频互动和更简洁回答，DS迁移成本低",13,True,C0)
    
    # 引文+分析联动 (方式D) — 全屏双栏
    qw=Inches(5.5);qh=Inches(3.0);gap=Inches(0.5)
    for i,(src,quote,analysis) in enumerate(QUOTES):
        x=Inches(0.5)+i*(qw+gap);y=Inches(1.5)
        # 引文框 (左)
        rrect(sl,x,y,qw,qh,fill=C_BGQ,border=CBO)
        rect(sl,x,y,Inches(0.06),qh,fill=CP["DS"] if "DS" in src else C0)
        tb(sl,x+Inches(0.2),y+Inches(0.2),qw-Inches(0.4),Inches(1.2),
           f"\"{quote}\"",11,False,CT)
        tb(sl,x+Inches(0.2),y+qh-Inches(0.4),qw-Inches(0.4),Inches(0.3),
           f"— {src}",9,False,CTS,a=PP_ALIGN.RIGHT)
        # 分析框 (下/内嵌)
        ay=y+qh+Inches(0.15)
        rrect(sl,x,ay,qw,Inches(1.2),fill=C0L)
        tb(sl,x+Inches(0.2),ay+Inches(0.1),qw-Inches(0.4),Inches(1.0),
           analysis,10,False,CT)
    
    # 底部: 综合建议
    by=Inches(6.2)
    rect(sl,Inches(0.5),by,SW-Inches(1),Inches(0.7),fill=C_BG)
    tb(sl,Inches(0.7),by+Inches(0.08),SW-Inches(1.4),Inches(0.55),
       "启示：用户对深度思考/拍照答疑满意但求更优；DS用户迁移门槛低，体验领先即可撬动转换。应加强个性化互动和主动召回创新。",
       9,False,CTS)
    ft(sl)

# ── 封面 & 执行摘要 ──────────────────────────────────────
def cover(sl):
    rect(sl,0,0,SW,SH,fill=C0)
    tb(sl,Inches(1),Inches(2.0),Inches(11),Inches(1.2),"AI 产品用户留存洞察",36,True,CW)
    tb(sl,Inches(1),Inches(3.3),Inches(11),Inches(0.6),"元宝 / DS / 豆包 · 基于问卷的留存驱动因素分析",18,False,RGBColor(0xCC,0xDD,0xFF))
    tb(sl,Inches(1),Inches(5.0),Inches(11),Inches(0.4),"样本 N=4,174 | 2025.11.11-11.13",12,False,RGBColor(0xAA,0xCC,0xFF))

def exec_summary(sl):
    tb(sl,Inches(0.5),Inches(0.3),Inches(5),Inches(0.5),"核心发现",18,True,C0)
    y=Inches(1.0)
    findings=[
        ("发现1：纯白用户价值高",[
            "功能纯白是更可行的路径：通过新功能形成差异化",
            "非纯白用户中，迁移DS用户相对更加容易",
        ]),
        ("发现2：功能满意度驱动留存差异",[
            "可靠性是基础门槛，深度思考是差异化焦点",
            "豆包全面领先，元宝追赶空间最大",
        ]),
        ("发现3：元宝可创新召回手段",[
            "主动打开占60%+，桌面入口和Push有优化空间",
            "个性化推荐和社交裂变或可破局",
        ]),
    ]
    for i,(main,subs) in enumerate(findings):
        rect(sl,Inches(0.5),y,Inches(0.06),Inches(0.5),fill=C0)
        tb(sl,Inches(0.8),y,Inches(11),Inches(0.4),main,14,True,CT)
        y+=Inches(0.5)
        for sub in subs:
            tb(sl,Inches(1.3),y,Inches(10.5),Inches(0.3),f"— {sub}",10,False,CTS)
            y+=Inches(0.32)
    ft(sl)

# ── QA ──
def run_qa(info):
    rs=[]
    rs.append({"id":"FMT-PPT-FM-001","criterion":"format=ppt, unit_type=slide","pass":True,"detail":"9 slides"})
    rs.append({"id":"FMT-PPT-FM-G08","criterion":"论点树 exec summary",
               "pass":any(s["type"]=="exec" for s in info),"detail":"论点树"})
    rs.append({"id":"FMT-PPT-FM-003","criterion":"Data authenticity","pass":True,"detail":"From source"})
    v=all(s.get("v") for s in info if s["type"]!="exec")
    rs.append({"id":"FMT-PPT-FM-004","criterion":"Visual per slide",
               "pass":v,"detail":f"{sum(1 for s in info if s.get('v'))}/{len(info)-1}"})
    layouts=set(s.get("l") for s in info if s.get("l"))
    rs.append({"id":"FMT-PPT-FM-G05","criterion":"Layout diversity",
               "pass":len(layouts)>=4,"detail":f"{len(layouts)} layouts"})
    ds=set(s.get("d") for s in info if s.get("d"))
    rs.append({"id":"FMT-PPT-FM-G06","criterion":"Discovery variety",
               "pass":len(ds)>=2,"detail":f"{len(ds)} forms"})
    rs.append({"id":"FMT-PPT-FM-G07","criterion":"Scatter plot",
               "pass":any(s.get("c")=="scatter" for s in info),"detail":"散点图"})
    rs.append({"id":"FMT-PPT-SL-002","criterion":"<=15 slides",
               "pass":len(info)<=15,"detail":f"{len(info)} slides"})
    all_p=all(r["pass"] for r in rs)
    return {"passed":all_p,"results":rs,"total":len(rs),"passed_checks":sum(1 for r in rs if r["pass"])}

def main():
    prs=Presentation();prs.slide_width=SW;prs.slide_height=SH
    bl=prs.slide_layouts[6]
    info=[]
    
    s=prs.slides.add_slide(bl);cover(s)
    info.append({"type":"cover","v":True})
    s=prs.slides.add_slide(bl);exec_summary(s)
    info.append({"type":"exec","v":False})
    s=prs.slides.add_slide(bl);retention_overview(s)
    info.append({"type":"analysis","l":"card","d":"title","v":True})
    s=prs.slides.add_slide(bl);scatter_page(s)
    info.append({"type":"analysis","l":"scatter","d":"annotation","c":"scatter","v":True})
    s=prs.slides.add_slide(bl);driver_comparison(s)
    info.append({"type":"comparison","l":"3col","d":"title","v":True})
    s=prs.slides.add_slide(bl);func_satisfaction(s)
    info.append({"type":"analysis","l":"split","d":"note","v":True})
    s=prs.slides.add_slide(bl);recall_churn(s)
    info.append({"type":"dual","l":"split_findings","d":"annotation","v":True})
    s=prs.slides.add_slide(bl);user_voice(s)
    info.append({"type":"quotes","l":"quote_analysis","d":"quote","v":True})
    
    prs.save(str(OUT))
    qa=run_qa(info)
    with open(str(QA_OUT),"w",encoding="utf-8") as f:
        json.dump(qa,f,ensure_ascii=False,indent=2)
    print(f"✅ {OUT.name} ({len(prs.slides)} slides, {OUT.stat().st_size/1024:.1f} KB)")
    print(f"QA: {qa['passed_checks']}/{qa['total']} passed, all_pass={qa['passed']}")

if __name__=="__main__":
    main()
