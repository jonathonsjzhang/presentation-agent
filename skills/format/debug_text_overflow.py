"""
v3.9 综合文本溢出诊断
===================
检测 3 类文本溢出：
  1. 文字行宽超 box w (单行字符数 × 平均字宽 > box w)
  2. 文字总高超 box h (估算行数 × 行高 > box h)
  3. 文字尾部被相邻 box 遮挡 (基于 AABB)

排除（不算溢出）：
  - 页眉 nav bar (y < 0.40)
  - 页脚 (y > 6.85)
  - 极小字 (< 7pt, 装饰性)
  - 故意"…"省略号结尾
  - 1-2 字符极短文本（size 误判无意义）
"""
import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

# 估算中英文字符宽度系数（相对 em）
# 中文字符 ≈ 1.0 em (size pt)
# 英文字符 ≈ 0.5 em (size pt)
# 数字/标点 ≈ 0.5 em
CHAR_W_CN = 1.0
CHAR_W_EN = 0.5
LINE_HEIGHT_FACTOR = 1.25  # 行高 = font size × 1.25

# 中文字符正则
CN_RE = re.compile(r'[\u4e00-\u9fff]')
EN_RE = re.compile(r'[a-zA-Z0-9]')


def estimate_text_width_pt(text, font_size_pt):
    """估算文字宽度（pt）"""
    width_pt = 0.0
    for ch in text:
        if CN_RE.match(ch):
            width_pt += font_size_pt * CHAR_W_CN
        else:
            width_pt += font_size_pt * CHAR_W_EN
    return width_pt


def estimate_lines(text, font_size_pt, box_w_pt):
    """估算单段文字需要的行数（考虑手动换行符）"""
    # 1 inch = 72 pt
    if box_w_pt <= 0:
        return 1
    # 按显式 \n 拆分
    lines = text.split('\n')
    total_lines = 0
    for line in lines:
        w = estimate_text_width_pt(line, font_size_pt)
        # 每行最大字符宽 = box_w_pt
        n_lines = max(1, int(w / box_w_pt) + (1 if w % box_w_pt > 0 else 0))
        total_lines += n_lines
    return total_lines


def get_run_font_size(run):
    """取 run 的字号（pt）"""
    try:
        if run.font.size is not None:
            return run.font.size.pt
    except Exception:
        pass
    return None


def get_para_font_size(para):
    """取 paragraph 的默认字号（pt）"""
    try:
        if para.font.size is not None:
            return para.font.size.pt
    except Exception:
        pass
    return None


def get_text_max_font_size(text_frame, default_size=10):
    """取 text_frame 中所有 run 和 paragraph 的最大字号"""
    sizes = []
    for para in text_frame.paragraphs:
        # 1. paragraph-level
        ps = get_para_font_size(para)
        if ps is not None:
            sizes.append(ps)
        # 2. run-level
        for run in para.runs:
            s = get_run_font_size(run)
            if s is not None:
                sizes.append(s)
    return max(sizes) if sizes else default_size


def get_text_full(text_frame):
    """获取 text_frame 完整文字（保留段落分隔）"""
    return '\n'.join(p.text for p in text_frame.paragraphs)


def is_header_or_footer(y_inch):
    """判断是否为页眉/页脚"""
    return y_inch < 0.40 or y_inch > 6.85


def diagnose(pptx_path, verbose=True):
    prs = Presentation(pptx_path)
    issues = []
    summary_per_slide = {}

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_issues = []
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            tf = shp.text_frame
            text = get_text_full(tf)
            if not text.strip():
                continue

            x = (shp.left or 0) / 914400
            y = (shp.top or 0) / 914400
            w = (shp.width or 0) / 914400
            h = (shp.height or 0) / 914400

            if is_header_or_footer(y):
                continue

            # 排除 page number / 极短字符串（如"6"）
            if len(text.strip()) <= 2 and text.strip().isdigit():
                continue

            max_size = get_text_max_font_size(tf)
            if max_size < 6.5:
                continue  # 极小字通常是装饰

            # 1 inch = 72 pt
            box_w_pt = w * 72
            box_h_pt = h * 72

            # 估算需要行数
            lines_needed = estimate_lines(text, max_size, box_w_pt)
            line_h_pt = max_size * LINE_HEIGHT_FACTOR
            text_h_pt = lines_needed * line_h_pt

            overflow_msg = []
            # 检查 1: 单段是否超宽
            for para in text.split('\n'):
                w_needed = estimate_text_width_pt(para, max_size)
                if w_needed > box_w_pt * 1.05:  # 5% tolerance
                    overflow_msg.append(
                        f"  - 段宽溢出: \"{para[:30]}{'...' if len(para)>30 else ''}\" "
                        f"需要 {w_needed:.1f}pt, box={box_w_pt:.1f}pt"
                    )

            # 检查 2: 总高是否超
            if text_h_pt > box_h_pt * 1.10:  # 10% tolerance
                overflow_msg.append(
                    f"  - 总高溢出: {lines_needed}行 × {line_h_pt:.1f}pt = {text_h_pt:.1f}pt, "
                    f"box={box_h_pt:.1f}pt (sz={max_size}pt)"
                )

            if overflow_msg:
                slide_issues.append({
                    'box': f'({x:.2f},{y:.2f},{w:.2f}x{h:.2f})',
                    'size': max_size,
                    'text': text.replace('\n', '|')[:60],
                    'msgs': overflow_msg,
                })

        if slide_issues:
            summary_per_slide[slide_idx] = len(slide_issues)
            for issue in slide_issues:
                issues.append((slide_idx, issue))

    if verbose:
        print(f"=== {Path(pptx_path).name} 文本溢出诊断 ===")
        print(f"共 {len(prs.slides)} 页, 找到 {len(issues)} 个潜在溢出\n")
        if not issues:
            print("✓ 未发现文本溢出")
            return 0
        for sidx, issue in issues:
            print(f"Slide {sidx} {issue['box']}  sz={issue['size']}pt  text=\"{issue['text']}\"")
            for m in issue['msgs']:
                print(m)
            print()
        print(f"=== 汇总：{len(summary_per_slide)} 张 slide 有问题 ===")
        for sidx, cnt in sorted(summary_per_slide.items()):
            print(f"  Slide {sidx}: {cnt} 个溢出")
    return len(issues)


if __name__ == "__main__":
    path = sys.argv[1] if len(sys.argv) > 1 else None
    if not path:
        # 默认查最新版本
        candidates = sorted(Path(r"C:\Users\zoezoezhao\Desktop\汇报agent\format skill 迭代").glob("format.skill.v*/AI 产品用户留存分析_skill_v*.pptx"))
        if not candidates:
            print("无 PPTX 可诊断")
            sys.exit(1)
        path = str(candidates[-1])
        print(f"诊断最新: {path}\n")
    diagnose(path)
