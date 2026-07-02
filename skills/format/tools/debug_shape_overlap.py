"""诊断 v3.6 PPT 的所有视觉重叠问题
- text-text: 文字-文字碰撞
- text-shape: 文字溢出 shape 边界
- shape-shape: shape 之间真实重叠（含图表 vs 文字框）
- chart-overlap: 图表 vs 文字框
"""
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX = r"C:\Users\zoezoezhao\Desktop\汇报agent\format skill 迭代\format.skill.v3.7\AI 产品用户留存分析_skill_v3.7.pptx"


def emu_to_in(emu):
    if emu is None:
        return 0.0
    return emu / 914400.0


def get_all_shapes(slide):
    """返回所有有 bbox 的 shape (含 chart)"""
    shapes = []
    for sh in slide.shapes:
        try:
            x = emu_to_in(sh.left)
            y = emu_to_in(sh.top)
            w = emu_to_in(sh.width)
            h = emu_to_in(sh.height)
        except Exception:
            continue
        if w < 0.02 or h < 0.02:
            continue
        is_chart = sh.has_chart
        is_text = sh.has_text_frame
        text = ""
        if is_text:
            text = " ".join(p.text for p in sh.text_frame.paragraphs).strip()[:50]
        shapes.append({
            "x": x, "y": y, "w": w, "h": h,
            "shape": sh,
            "is_chart": is_chart,
            "is_text": is_text,
            "text": text,
            "shape_id": sh.shape_id,
            "name": sh.name,
        })
    return shapes


def aabb(a, b):
    return (
        a["x"] < b["x"] + b["w"] and
        a["x"] + a["w"] > b["x"] and
        a["y"] < b["y"] + b["h"] and
        a["y"] + a["h"] > b["y"]
    )


def overlap_area(a, b):
    ox = max(0, min(a["x"] + a["w"], b["x"] + b["w"]) - max(a["x"], b["x"]))
    oy = max(0, min(a["y"] + a["h"], b["y"] + b["h"]) - max(a["y"], b["y"]))
    return ox * oy


def is_container(s):
    """判断 shape 是否为容器（panel 圆角矩形 / 背景框）"""
    if s["w"] < 0.5 or s["h"] < 0.4:
        return False
    return not s["is_text"] or not s["text"]


def is_header_footer_strip(s):
    """判断是否为顶/底横条（页眉/页脚）"""
    if s["w"] > 8 and s["h"] < 0.7:
        if s["y"] < 0.6 or s["y"] + s["h"] > 6.5:
            return True
    return False


def is_contained(inner, outer, pad=0.05):
    """inner 是否完全在 outer 内（容差 pad）"""
    return (inner["x"] >= outer["x"] - pad and
            inner["y"] >= outer["y"] - pad and
            inner["x"] + inner["w"] <= outer["x"] + outer["w"] + pad and
            inner["y"] + inner["h"] <= outer["y"] + outer["h"] + pad)


def text_overflow_check(slide, shapes):
    """检查文字是否溢出其 shape 边界"""
    issues = []
    for s in shapes:
        if not s["is_text"] or not s["text"]:
            continue
        # 估算文字宽度（简单估算: size_pt * 0.55 * char_count / 72 inch）
        # 取最大字符的 paragraph
        for p in s["shape"].text_frame.paragraphs:
            if not p.runs:
                continue
            for r in p.runs:
                if not r.text or not r.font.size:
                    continue
                # 估算宽高
                char_w_in = r.font.size.pt * 0.55 / 72 * len(r.text) * 0.6  # 经验系数
                line_h_in = r.font.size.pt * 1.3 / 72
                # 段落中所有 run
                pass
        # 简化: 用 text length * size 估算
        if not s["shape"].text_frame.paragraphs:
            continue
        max_run = max(
            ((r, p) for p in s["shape"].text_frame.paragraphs for r in p.runs if r.text and r.font.size),
            key=lambda rp: len(rp[0].text) * rp[0].font.size.pt if rp[0].text and rp[0].font.size else 0,
            default=(None, None)
        )
        if max_run[0] is None or not max_run[0].font.size:
            continue
        run = max_run[0]
        est_w = run.font.size.pt * 0.13 * len(run.text)  # 中文约 0.13 in/pt * char
        # 对中文 + 英文混合做修正
        chinese_chars = sum(1 for c in run.text if '\u4e00' <= c <= '\u9fff')
        english_chars = len(run.text) - chinese_chars
        est_w_zh = chinese_chars * run.font.size.pt * 0.16 / 72
        est_w_en = english_chars * run.font.size.pt * 0.07 / 72
        est_w = est_w_zh + est_w_en
        if est_w > s["w"] * 1.1 and est_w > 0.3:  # 溢出 10%+
            issues.append({
                "type": "text_overflow",
                "shape": s,
                "est_w": est_w,
                "box_w": s["w"],
            })
    return issues


def main():
    prs = Presentation(PPTX)
    all_issues = {}

    for idx, slide in enumerate(prs.slides, 1):
        shapes = get_all_shapes(slide)
        slide_issues = []

        # 1. text overflow check
        slide_issues.extend(text_overflow_check(slide, shapes))

        # 2. shape-shape overlap (parent-child excluded, header/footer excluded)
        for i, a in enumerate(shapes):
            for b in shapes[i+1:]:
                if not aabb(a, b):
                    continue
                area = overlap_area(a, b)
                if area < 0.02:  # 微小重叠忽略
                    continue
                # 排除父子关系
                if is_contained(a, b) or is_contained(b, a):
                    continue
                # 排除页眉/页脚
                if is_header_footer_strip(a) and is_header_footer_strip(b):
                    continue
                # 排除文字-shape 容器
                if a["is_text"] and not b["is_text"] and not b["is_chart"]:
                    # text inside shape (panel)?
                    if is_contained(a, b, pad=0.15):
                        continue
                if b["is_text"] and not a["is_text"] and not a["is_chart"]:
                    if is_contained(b, a, pad=0.15):
                        continue
                # 报告
                slide_issues.append({
                    "type": "shape_overlap",
                    "a": a, "b": b,
                    "area": area,
                })

        if slide_issues:
            all_issues[idx] = slide_issues

    # 输出
    if not all_issues:
        print("✓ 未发现 shape-shape 重叠 / 文字溢出")
        return

    print(f"发现 {len(all_issues)} 个 slide 有问题:\n")
    for idx, issues in all_issues.items():
        print(f"=== Slide {idx} ({len(issues)} issues) ===")
        for iss in issues:
            if iss["type"] == "text_overflow":
                s = iss["shape"]
                print(f"  [TEXT_OVERFLOW] ({s['x']:.2f},{s['y']:.2f},{s['w']:.2f}x{s['h']:.2f}) '{s['text'][:30]}' 估计宽={iss['est_w']:.2f} > 容器宽={iss['box_w']:.2f}")
            elif iss["type"] == "shape_overlap":
                a, b = iss["a"], iss["b"]
                a_kind = "chart" if a["is_chart"] else ("text" if a["is_text"] else "shape")
                b_kind = "chart" if b["is_chart"] else ("text" if b["is_text"] else "shape")
                a_text = a['text'][:25] if a['is_text'] else a['name']
                b_text = b['text'][:25] if b['is_text'] else b['name']
                print(f"  [SHAPE_OVERLAP area={iss['area']:.3f}]")
                print(f"    A [{a_kind}] ({a['x']:.2f},{a['y']:.2f},{a['w']:.2f}x{a['h']:.2f}) {a_text}")
                print(f"    B [{b_kind}] ({b['x']:.2f},{b['y']:.2f},{b['w']:.2f}x{b['h']:.2f}) {b_text}")
        print()


if __name__ == "__main__":
    main()
