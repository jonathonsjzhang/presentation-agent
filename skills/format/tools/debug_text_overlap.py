"""诊断 v3.2 PPT 的真实文字-文字重叠问题
策略: 忽略所有空 text frame 和被大容器完全包含的小 text（设计嵌入）
      只关注内容 text frame 之间的字符级碰撞
"""
from pptx import Presentation
from pptx.util import Emu

PPTX = r"C:\Users\zoezoezhao\Desktop\汇报agent\format skill 迭代\format.skill.v3.5\AI 产品用户留存分析_skill_v3.5.pptx"

def get_text_shapes(slide):
    """返回所有有实际文字内容的 text shape"""
    shapes = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if not sh.left or not sh.top or not sh.width or not sh.height:
            continue
        tf = sh.text_frame
        full_text = " ".join(p.text for p in tf.paragraphs).strip()
        if not full_text:
            continue
        # 跳过页码之类的微小元素
        x = sh.left / 914400
        y = sh.top / 914400
        w = sh.width / 914400
        h = sh.height / 914400
        if w * h < 0.05:  # 太小
            continue
        shapes.append({
            "x": x, "y": y, "w": w, "h": h,
            "text": full_text,
            "raw": sh,
        })
    return shapes

def is_contained(inner, outer, pad=0.05):
    """inner 是否完全在 outer 内（容差 pad）"""
    return (inner["x"] >= outer["x"] - pad and
            inner["y"] >= outer["y"] - pad and
            inner["x"] + inner["w"] <= outer["x"] + outer["w"] + pad and
            inner["y"] + inner["h"] <= outer["y"] + outer["h"] + pad)

def is_subtitle_or_strip(inner, outer, pad=0.10):
    """inner 是否在 outer 的正下方或侧边（看作标题/页脚条带）"""
    # 顶部条带 (subtitle, top frame)
    top_strip = (inner["y"] < 1.2 and outer["y"] < 1.2 and
                 inner["w"] > 5 and outer["w"] > 5)
    # 底部条带 (footer, data strip)
    bot_strip = (inner["y"] > 5.5 and outer["y"] > 5.5 and
                 inner["w"] > 5 and outer["w"] > 5)
    return top_strip or bot_strip

prs = Presentation(PPTX)
report = []

for idx, slide in enumerate(prs.slides, 1):
    # 跳过封面/封底（无要求统一边距）
    if idx in (1, 16):
        continue
    text_shapes = get_text_shapes(slide)
    issues = []
    for i, t1 in enumerate(text_shapes):
        for t2 in text_shapes[i+1:]:
            # AABB 重叠
            aabb_overlap = (
                t1["x"] < t2["x"] + t2["w"] and
                t1["x"] + t1["w"] > t2["x"] and
                t1["y"] < t2["y"] + t2["h"] and
                t1["y"] + t1["h"] > t2["y"]
            )
            if not aabb_overlap:
                continue
            # 排除"容器-子元素"嵌套
            t1_in_t2 = is_contained(t1, t2)
            t2_in_t1 = is_contained(t2, t1)
            if t1_in_t2 or t2_in_t1:
                continue
            # 排除"标题/页脚条带"互撞（设计就是水平条带）
            if is_subtitle_or_strip(t1, t2):
                continue
            # 真实文字-文字碰撞
            issues.append(f"  ({t1['x']:.2f},{t1['y']:.2f},{t1['w']:.2f}x{t1['h']:.2f}) '{t1['text'][:30]}'")
            issues.append(f"  ×")
            issues.append(f"  ({t2['x']:.2f},{t2['y']:.2f},{t2['w']:.2f}x{t2['h']:.2f}) '{t2['text'][:30]}'")
            issues.append("")
    if issues:
        report.append(f"=== Slide {idx} ===")
        report.extend(issues)

if report:
    print(f"\n发现 {sum(1 for r in report if r.startswith('==='))} 个 slide 有真实文字重叠:\n")
    print("\n".join(report))
else:
    print("✓ 未发现真实文字-文字重叠")
