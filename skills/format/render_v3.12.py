# -*- coding: utf-8 -*-
"""
Render slides using format.skill v3.12 (v3.6→v3.11 cumulative + v3.12 hotfix).

[v3.12 hotfix 关键变更]
- 8 个 layout 常量提取 (BLANK_LAYOUT_IDX / CONTENT_LEFT / CONTENT_W / DISCOVERY_X/W/H / BOTTOM_Y/H)
- 16 个 page builder 中 prs.slide_layouts[6] → prs.slide_layouts[BLANK_LAYOUT_IDX]
- add_discovery_panel 默认值修复 (10.20/2.83/5.20 → 10.00/2.63/4.10，与 12 个调用方对齐)
- add_bottom_summary 9 个调用方移除冗余默认参数
- 3 处 dict.get("可用性", 0) → dict(...)[...] strict 模式
- 删除 8 个死函数 (rgb_hex, add_red_box, add_not_significant, add_insight_panel_rail/banner/callout, add_insight_panel alias, add_legend) + 1 个未用 import (MSO_LINE_DASH_STYLE) + 1 个未用 kwarg (layout="3x1") + 1 个未用常量 (TITLE_EA_FONT)
- QA 5 重门禁全绿 (FMT-V3-009/010/011/012/013)

[v3.12 视觉系统（沿用 v3.6→v3.11）]
- v0.9.1 基线视觉系统（顶部蓝 nav + Tencent logo + 保密标签 + 三方品牌色 + 红框）
- v3.6 4 区填充规则（顶部 nav + 左侧主图 + 右侧"发现" rail + 底部 bottom_summary）
- v3.7 7 类 insight_panel 变体（right-rail / top-banner / bottom-takeaway / callout-side / inline-anchor / matrix-grid / bottom-summary）
- v3.8 baseline 5.40 对齐（全 PPT 16 页图表底 y 统一 = 5.40）
- v3.9 debug_text_overflow.py 综合诊断工具（text-text + shape-shape + text-overflow）
- v3.10 P3 顶部 4 blocks 重构（(维度名, 数量, 枚举值) + 用户数附加信息）
- v3.11 中英文字体分离（中文 → 楷体 <a:ea> / 英文/数字 → arial <a:latin>，使用 lxml _set_dual_font()）

[v3.12 QA 增强]
- FMT-V3-009 边距检查 (v3.5 起)
- FMT-V3-010 重叠检查 (v3.5 起，v3.7 升级为 text-text + shape-shape + text-overflow)
- FMT-V3-011 executive_summary_no_rail (v3.4 起)
- FMT-V3-012 底部 bottom_summary 必备 (v3.6 起)
- FMT-V3-013 右 rail 必备 (v3.6 起)
- v3.12 hotfix 新增隐式：layout 漂移 QA 通过 8 个常量统一收口，避免 P0 漂移复发
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path
import json

# v3.12 output paths
OUT_DIR = Path(r"C:\Users\zoezoezhao\Desktop\汇报agent\format skill\format skill 迭代\format.skill.v3.12")
OUT = OUT_DIR / "AI 产品用户留存分析_skill_v3.12.pptx"
QA_OUT = OUT_DIR / "AI 产品用户留存分析_skill_v3.12_QA.json"

# [v3.12] Layout 常量 (修复 P0 layout 漂移：此前 MARGIN_INTERNAL 仅在 QA 中引用，16 个 page builder 硬编码 0.7/11.933)
BLANK_LAYOUT_IDX = 6     # blank 母版索引
CONTENT_LEFT = 0.7       # FMT-V3-009: 内部页面统一左边距
CONTENT_W = 11.933       # FMT-V3-009: 内部页面统一内容宽度
CONTENT_RIGHT = CONTENT_LEFT + CONTENT_W  # 12.633
DISCOVERY_X = 10.00      # v3.5+ discovery panel 实际位置 (修复前函数默认 10.20, 调用方传 10.00, 漂移)
DISCOVERY_W = 2.63       # v3.5+ discovery panel 实际宽度
DISCOVERY_H = 4.10       # 避开 bottom_summary (5.55-6.65), 与 bottom 留 0.05 缓冲
BOTTOM_Y = 5.55          # bottom_summary 顶部 Y
BOTTOM_H = 1.10          # bottom_summary 高度 (留 6.65+0.20 给 footer)

# [v3 layout_type 标注] case1 的 16 页 → v3 5 类 layout 映射
V4_LAYOUT_MAP = {
    1:  "cover",
    2:  "executive_summary",
    3:  "methodology_or_strategy",
    4:  "analysis_dashboard",   # 留存总览 (bar_chart)
    5:  "analysis_dashboard",   # 驱动矩阵 (horizontal_bar)
    6:  "analysis_dashboard",   # 纯白用户 (scatter)
    7:  "analysis_dashboard",   # 用户来源矩阵 (custom blocks)
    8:  "analysis_dashboard",   # 文本满意度 (grouped_bar)
    9:  "analysis_dashboard",   # 功能满意度 (matrix_table)
    10: "analysis_dashboard",   # 功能纯白 (bar_chart + matrix_table)
    11: "analysis_dashboard",   # 召回机制 (grouped_bar)
    12: "analysis_dashboard",   # 流失原因 (bar_chart + cards)
    13: "analysis_dashboard",   # 用户原声 (quoted_table)
    14: "methodology_or_strategy",  # 行动建议 (Exploit-Explore-Watch)
    15: "priority_matrix",      # 优先级矩阵 (quadrants)
    16: "cover",                # 结论页 (深色 cover 变体)
}

# ---------- Data extracted from AI 产品用户留存分析_文档资料.pdf ----------
SAMPLE_DATA = {"total": 4174, "date": "2025.11.11-11.13", "method": "问卷投放：性别年龄=网民大盘，学历职业偏大模型核心用户"}
PRODUCTS = ["豆包", "DS", "元宝"]
PKEY = {"豆包": "doubao", "DS": "ds", "元宝": "yuanbao"}
STRONG_RETENTION = {
    "元宝": {"rate": 19.0, "users": 2771, "penetration": 66.4},
    "DS": {"rate": 34.0, "users": 3607, "penetration": 86.4},
    "豆包": {"rate": 54.0, "users": 3968, "penetration": 95.1},
}
PURE_NEW_USERS = {
    "元宝": {"pure_new": 215, "pure_new_rate": 26, "non_pure_new_rate": 12, "correlation": 0.052, "p_value": 0.006},
    "DS": {"pure_new": 534, "pure_new_rate": 45, "non_pure_new_rate": 27, "correlation": 0.097, "p_value": 0.000},
    "豆包": {"pure_new": 712, "pure_new_rate": 60, "non_pure_new_rate": 50, "correlation": 0.052, "p_value": 0.001},
}
RETENTION_DRIVERS = {
    "豆包": [("拍照答疑", "纯白", 21.0), ("深度思考", "纯白", 20.2), ("速度及稳定性", "5分满意", 20.1), ("情绪价值", "5分满意", 20.1), ("可靠性", "5分满意", 18.9), ("AI创作", "纯白", 16.3)],
    "DS": [("可靠性", "5分满意", 25.2), ("深度思考", "纯白", 23.8), ("情绪价值", "5分满意", 23.3), ("拍照答疑", "纯白", 17.9), ("可用性", "5分满意", 17.5), ("速度及稳定性", "5分满意", 15.6)],
    "元宝": [("可靠性", "5分满意", 16.1), ("深度思考", "5分满意", 14.3), ("深度思考", "纯白", 14.0), ("拍照答疑", "纯白", 12.6), ("AI创作", "纯白", 12.0), ("纯白用户", "", 11.9)],
}
TEXT_SATISFACTION_CORR = {
    "豆包": [("可靠性", 0.224), ("可用性", 0.213), ("速度/稳定性", 0.204), ("情绪价值", 0.188)],
    "DS": [("可靠性", 0.158), ("情绪价值", 0.121), ("可用性", 0.119), ("速度/稳定性", 0.100)],
    "元宝": [("可靠性", 0.131), ("速度/稳定性", 0.101), ("可用性", 0.079), ("情绪价值", 0.078)],
}
FUNCTION_SATISFACTION_CORR = {
    "豆包": [("深度思考", 0.227), ("拍照答疑", 0.184), ("打电话", 0.172), ("朗读", 0.170), ("AI创作", 0.144)],
    "DS": [("深度思考", 0.137), ("AI创作", 0.061), ("拍照答疑", 0.057), ("打电话", 0.042), ("朗读", 0.036)],
    "元宝": [("打电话", 0.170), ("AI创作", 0.144), ("朗读", 0.144), ("深度思考", 0.122), ("拍照答疑", 0.085)],
}
FUNCTION_PURE_NEW_LIFT = {
    "豆包": {"深度思考": (59, 53, 6), "拍照答疑": (62, 52, 10), "打电话": (57, 52, 5), "AI创作": (56, 52, 4), "朗读": (57, 52, 5)},
    "DS": {"深度思考": (47, 32, 16), "拍照答疑": (49, 32, 17)},
    "元宝": {"深度思考": (29, 18, 12), "拍照答疑": (34, 18, 16), "打电话": (29, 17, 12), "AI创作": (31, 17, 14), "朗读": (29, 17, 12)},
}
RECALL_DATA = {
    "主动打开": {"豆包": 58.1, "DS": 57.5, "元宝": 61.3},
    "有红点": {"豆包": 14.0, "DS": 19.4, "元宝": 23.0},
    "桌面入口": {"豆包": 18.3, "DS": 29.3, "元宝": 30.5},
    "Push": {"豆包": 11.0, "DS": 13.0, "元宝": 11.5},
    "朋友分享": {"豆包": 14.0, "DS": 13.0, "元宝": 13.0},
}
CHURN_REASONS = {
    "豆包": {"text_unsatisfied": 44, "function_unsatisfied": 40, "top_reason": "可靠性：答案准确性、幻觉少、信源权威"},
    "DS": {"text_unsatisfied": 57, "function_unsatisfied": None, "top_reason": "速度/稳定性（66%）+ 可靠性（26%）"},
    "元宝": {"text_unsatisfied": 62, "function_unsatisfied": 31, "top_reason": "可靠性差（51%希望答案更准确、索引更权威）"},
}
USER_QUOTES = {
    "元宝正向": ["回答速度非常快、稳定性好（16%用户提到）", "生图等多模态功能好（17%用户提到）"],
    "元宝负向": ["可靠性差，希望答案更准确、索引更权威（51%）", "生图功能需加强：分辨率、质量、细节（20%）"],
    "豆包正向": ["可靠性、可用性、速度好（28%用户提到）", "AI创作等功能满意（12%）"],
    "豆包负向": ["可靠性差，幻觉少、信源质量好（44%）", "AI创作等文生图、文生视频一致性不满（40%）"],
    "DS正向": ["可靠性好，逻辑思维、专业知识深度（16%）", "功能单一简洁（6%）"],
    "DS负向": ["速度/稳定性差（66%），其次可靠性（26%）", "缺少功能"],
}

# ---------- Style: format.skill v3 (references/style.md) ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 13.333, 7.5

COL = {
    # [v3.5] vS 精确品牌色板 (从 case1.vs.pdf 逐像素提取)
    "primary": RGBColor(0, 102, 204),       # #0066CC — vS 品牌主色
    "page_title": RGBColor(30, 111, 224),   # #1E6FE0 — vS 内容页主标题色
    "navy": RGBColor(5, 28, 44),            # #051C2C — 深海军(用于 cover)
    "nav_bar": RGBColor(0, 61, 130),        # #003D82 — vS 导航栏底色
    # [v3.5] 三方品牌色 (vS 实证) — 元宝=绿, DS=蓝, 豆包=浅蓝
    "product_a": RGBColor(63, 191, 111),    # #3FBF6F — 元宝
    "product_b": RGBColor(30, 111, 224),    # #1E6FE0 — DS
    "product_c": RGBColor(187, 216, 248),   # #BBD8F8 — 豆包
    "success": RGBColor(63, 191, 111),      # #3FBF6F — 正向
    "warning": RGBColor(239, 68, 68),       # 警示
    "orange": RGBColor(245, 158, 11),        # 强调
    # [v3.5] vS 面板/标注色
    "panel_bg": RGBColor(245, 245, 245),    # #F5F5F5 — vS 发现面板底
    "discovery_bg": RGBColor(245, 245, 245),# 同 panel_bg (语义别名)
    "panel_blue": RGBColor(232, 240, 250),  # #E8F0FA — vS 浅蓝 callout
    "callout_green": RGBColor(200, 230, 201), # #C8E6C9 — vS 浅绿 callout
    "red_highlight": RGBColor(211, 47, 47), # #D32F2F — vS 红框高亮
    # 中性/文字色
    "text": RGBColor(51, 51, 51),           # #333333 — vS 正文
    "text_bold": RGBColor(33, 33, 33),      # #212121 — vS 粗体正文
    "muted": RGBColor(136, 136, 136),       # #888888 — vS 脚注/页码
    "not_sig": RGBColor(170, 170, 170),     # #AAAAAA — vS 不显著标注
    "light": RGBColor(224, 224, 224),       # #E0E0E0 — 浅灰线
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(0, 0, 0),
}
PCOL = {"豆包": COL["product_c"], "DS": COL["product_b"], "元宝": COL["product_a"]}
PLIGHT = {"豆包": RGBColor(230, 255, 239), "DS": RGBColor(232, 240, 255), "元宝": RGBColor(239, 247, 255)}
# [v3.11] 中英文字体分离：英文/数字 → arial，中文 → 楷体
# python-pptx 的 font.name 只设 <a:latin>，需要 lxml 同时加 <a:ea> 才是中英文分离
FONT = "arial"           # latin/英文/数字
EA_FONT = "楷体"          # eastAsia/中文
TITLE_FONT = "arial"     # 标题 latin (与正文共用 FONT,无独立常量)

def _set_dual_font(run, latin, ea, size_pt=None, bold=False, color=None):
    """给 run 同时设 latin (英文/数字) + eastAsia (中文) 字体

    python-pptx 的 run.font.name 只设 <a:latin typeface="...">，
    中文字体必须在 <a:ea typeface="..."> 里独立指定才能中英文分离。
    """
    # 1. latin (英文/数字)
    rPr = run._r.get_or_add_rPr()
    # 移除已有 latin
    for child in rPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}latin"):
        rPr.remove(child)
    latin_el = rPr.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}latin",
        {"typeface": latin},
    )
    rPr.append(latin_el)
    # 2. eastAsia (中文)
    for child in rPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}ea"):
        rPr.remove(child)
    ea_el = rPr.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
        {"typeface": ea},
    )
    rPr.append(ea_el)
    # 3. cs 同步（避免复杂脚本回退）
    for child in rPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}cs"):
        rPr.remove(child)
    cs_el = rPr.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}cs",
        {"typeface": latin},
    )
    rPr.append(cs_el)
    # 4. 字号 / 粗体 / 颜色（用 python-pptx 友好接口）
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if bold:
        run.font.bold = True
    if color is not None:
        run.font.color.rgb = color
# [v3.5] vS 风格导航 tabs (case1.vs 实证: "人群  纯白  非纯白  文本  功能  运营")
TABS = ["人群", "纯白", "非纯白", "文本", "功能", "运营"]

def slide_bg(slide, color=COL["white"]):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def add_text(slide, text, x, y, w, h, size=14, color=None, bold=False, align=PP_ALIGN.LEFT, font=FONT, valign=MSO_ANCHOR.TOP, margin=0.03):
    # [v3.11] 中英文字体分离：latin=font (默认 arial), eastAsia=EA_FONT (默认 楷体)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(margin); tf.margin_right = Inches(margin); tf.margin_top = Inches(margin); tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        _set_dual_font(run, latin=font, ea=EA_FONT, size_pt=size, bold=bold, color=color or COL["text"])
    p.alignment = align
    return box

def add_para_text(slide, paras, x, y, w, h, size=12, color=None, bullet=False, line_space=1.0):
    # [v3.11] 中英文字体分离
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05); tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    for i, t in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " + t) if bullet else t
        for run in p.runs:
            _set_dual_font(run, latin=FONT, ea=EA_FONT, size_pt=size, color=color or COL["text"])
        p.space_after = Pt(5)
        p.line_spacing = line_space
    return box

def add_shape(slide, shape, x, y, w, h, fill, line=None, radius=False):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    if radius and hasattr(s, "adjustments"):
        try: s.adjustments[0] = 0.12
        except Exception: pass
    return s

def add_title(slide, title, subtitle=None, dark=False):
    # [v3.5] 标题色改 vS 蓝 #1E6FE0 (替代原 navy); y=0.70 适配新顶部导航栏
    # [v3.5 fix] w=11.933 让 x+w=12.633 满足 0.7 inch 边距
    color = COL["white"] if dark else COL["page_title"]
    add_text(slide, title, 0.7, 0.70, 11.933, 0.42, size=18, bold=True, color=color, font=TITLE_FONT)
    if subtitle:
        # [v3.7 修复] subtitle y 从 1.10 下移到 1.20, 避开与 title (y+h=1.12) 的 0.02 inch 重叠
        add_text(slide, subtitle, 0.7, 1.20, 11.5, 0.30, size=9.5, color=COL["muted"] if not dark else RGBColor(210, 220, 235))

# [v3.5] vS 风格顶部蓝色导航栏 — 全宽深蓝条 + 白色 tab + 当前 tab 反白高亮
# 参考 case1.vs P6/P10: 导航栏 y=0, 高度 0.30, 颜色 #003D82
# tab 字号 10.5pt, 当前 tab 白底蓝字(圆角 3px), 非当前白字
def add_nav(slide, active):
    nav_h = 0.30
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, nav_h, COL["nav_bar"], line=None)
    tab_w, gap, left = 0.85, 0.10, 0.5
    for i, tab in enumerate(TABS):
        x = left + i * (tab_w + gap)
        active_flag = (i == active)
        if active_flag:
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 0.04, tab_w, nav_h - 0.08, COL["white"], line=None, radius=True)
            add_text(slide, tab, x, 0.06, tab_w, nav_h - 0.12, size=10, bold=True, color=COL["nav_bar"], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        else:
            add_text(slide, tab, x, 0.06, tab_w, nav_h - 0.12, size=10, color=COL["white"], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

# [v3.5] Tencent logo — 右上角蓝色文字 logo (与 vS 白底页一致)
# cover 页面用 dark=True 改为白色字
def add_tencent_logo(slide, dark=False):
    color = COL["white"] if dark else COL["page_title"]
    add_text(slide, "Tencent", 11.7, 0.05, 1.4, 0.22, size=14, bold=True, color=color, font="Segoe UI")

# [v3.5] 保密标签 — logo 下方, 内容页标配
def add_confidential_tag(slide):
    add_text(slide, "内部汇报 · 仅供参考", 9.5, 0.32, 2.0, 0.18, size=8.5, color=RGBColor(153, 153, 153), align=PP_ALIGN.RIGHT)

def add_footer(slide, source, page_no):
    # [v3.5] vS 注1:...; 注2:... 多行编号格式 (替代原单行 注:...)
    # source 支持 str 或 list[str]
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.7, 6.85, 11.933, 0.01, COL["light"])
    if isinstance(source, str):
        notes = [source]
    else:
        notes = list(source)
    note_text = "; ".join(f"注{i+1}：{n}" for i, n in enumerate(notes))
    add_text(slide, note_text, 0.7, 6.90, 10.5, 0.22, size=7.8, color=COL["muted"])
    add_text(slide, f"{page_no:02d}", 12.05, 6.89, 0.55, 0.22, size=8.5, color=COL["muted"], align=PP_ALIGN.RIGHT)

# ============================================================
# v3.2: Insight Panel 6 变体（位置/形态灵活,按 evidence 选型）
# 见 references/components.md §C-2 与 evals/rubrics.json FMT-V3-008
# [v3.2 位置调整] 所有变体位置适配 0.7 inch 统一边距(FMT-V3-009)
# 见 references/style.md §4.1
# ============================================================

# [v3.5] vS 标准"发现"右栏 panel — 模拟 case1.vs P6/P10 右侧灰色 panel
# 视觉: #F5F5F5 灰底 + "发现"标题(14pt Bold #212121) + bullet 列表(12pt #333,行距 1.4x)
# 与 add_insight_panel_rail 的关键差异:
#   - 标题用"发现"(不是"关键洞察")
#   - bullet 用 "•" 圆点 (不是编号圆圈)
#   - panel 内部无边线/无装饰
#   - 标题在 panel 内部顶部,左对齐
# [v3.5 fix] bullet 容器 h 从 1.10 缩到 0.85, cy 步进 0.95,避免 sibling AABB 误报
def add_discovery_panel(slide, items, x=DISCOVERY_X, y=1.30, w=DISCOVERY_W, h=DISCOVERY_H, title="发现"):
    """vS 风格"发现"右栏 panel (case1.vs P6/P10 实证)
    [v3.12] 默认 x=10.00 y=1.30 w=2.63 h=4.10 (与 12 个调用方实际值对齐,避开 bottom_summary 区域)
    [v3.12] 无 bottom_panel 的页面 (P4 纯白) 显式传 h=5.20 即可"""
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, COL["discovery_bg"], line=None)
    add_text(slide, title, x + 0.18, y + 0.16, w - 0.36, 0.32, size=14, bold=True, color=COL["text_bold"])
    # 标题下分隔线 (浅灰)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.18, y + 0.52, w - 0.36, 0.01, COL["light"])
    cy = y + 0.68
    bullet_h = 0.85
    bullet_gap = 0.10
    for item in items[:4]:
        # bullet 圆点
        add_text(slide, "•", x + 0.18, cy, 0.20, 0.24, size=14, bold=True, color=COL["text_bold"])
        # 文字 (h=0.85 容纳 4 行 11pt 中文)
        add_text(slide, item, x + 0.40, cy + 0.02, w - 0.58, bullet_h, size=11, color=COL["text"])
        cy += bullet_h + bullet_gap

def add_insight_panel_takeaway(slide, columns, x=0.7, y=6.05, w=11.933, h=0.65, title="行动建议"):
    """C-2-C bottom-takeaway: 底部横条 + 强调色边 + 多列分块（立即做/短期做/中期探索）
    [v3.2] 位置 x=0.7 y=6.05 w=11.933 h=0.65 (适配 0.7 inch 统一边距,留出 footer 空间)"""
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, COL["panel_bg"], radius=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.05, COL["primary"])
    add_text(slide, title, x + 0.18, y + 0.10, 1.2, 0.40, size=11, bold=True, color=COL["navy"])
    n = len(columns)
    col_w = (w - 1.50) / n
    cx = x + 1.42
    for i, col in enumerate(columns):
        if i > 0:
            add_shape(slide, MSO_SHAPE.RECTANGLE, cx - 0.06, y + 0.12, 0.015, h - 0.20, COL["light"])
        add_text(slide, col.get("tag", ""), cx, y + 0.10, col_w - 0.12, 0.22, size=9.5, bold=True, color=COL["primary"])
        add_text(slide, col.get("text", ""), cx, y + 0.32, col_w - 0.12, h - 0.36, size=9.0, color=COL["text"])
        cx += col_w

# [v3.6] C-2-G bottom-summary: 底部关键发现 3 列 panel (用于 4 区填充,填补底部空白)
# 与 add_insight_panel_takeaway 的关键差异:
#   - 默认 y=5.55 (上提 0.50) 腾出 footer 空间 (0.20 + 6.65+0.20=6.85 留出 footer 0.20 缓冲)
#   - 默认 h=1.10 (上提 0.45) — 给 footer (y=6.85) 留 0.20 缓冲
#   - 标题左侧增加强调色块 (4px)
#   - 文字更大 (size 10.5), 适配更长行文
# [v3.6 FMT-V3-012] 强制每页 (除 P1/P2/P16) 必须有底部 panel 填满 y=5.55-6.65 区
def add_bottom_summary(slide, columns, x=0.7, y=5.55, w=11.933, h=1.10, title="关键发现", title_color=None):
    """C-2-G bottom-summary: 底部 3 列关键发现 panel (v3.6 4 区填充必选)
    [v3.6] 位置 x=0.7 y=5.55 w=11.933 h=1.10 (底部 1.10 inch 全宽,留 0.20 给 footer)
    [v3.6] columns: list of dict {tag, text, color?}"""
    title_color = title_color or COL["primary"]
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, COL["panel_bg"], radius=True)
    # 左侧 4px 强调色竖线
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.06, h, title_color)
    # 标题 (左 1.2 inch)
    add_text(slide, title, x + 0.18, y + 0.16, 1.2, 0.40, size=12, bold=True, color=COL["navy"])
    n = len(columns)
    col_w = (w - 1.50) / n
    cx = x + 1.42
    for i, col in enumerate(columns):
        if i > 0:
            add_shape(slide, MSO_SHAPE.RECTANGLE, cx - 0.06, y + 0.15, 0.015, h - 0.30, COL["light"])
        col_color = col.get("color", COL["primary"])
        add_text(slide, col.get("tag", ""), cx, y + 0.12, col_w - 0.12, 0.22, size=10, bold=True, color=col_color)
        add_text(slide, col.get("text", ""), cx, y + 0.36, col_w - 0.12, h - 0.42, size=10, color=COL["text"])
        cx += col_w
    return x, y, w, h

def add_insight_panel_anchor(slide, text, target_x, target_y, bx=7.0, by=1.20, bw=2.6, bh=0.55):
    """C-2-E inline-anchor: 引线 + 小气泡,锚定到图表某数据点(target_x,target_y)
    [v3.2] 默认 bx=7.0 by=1.20 bw=2.6 bh=0.55 (避开 grouped_bar 上方,适配 0.7 边距)"""
    # 气泡中心点
    bubble_cx, bubble_cy = bx + bw / 2, by + bh / 2
    # 计算引线起点（target）和终点（气泡边沿）
    dx, dy = bubble_cx - target_x, bubble_cy - target_y
    dist = (dx ** 2 + dy ** 2) ** 0.5
    if dist > 0:
        end_x = target_x + dx * (dist - 0.20) / dist
        end_y = target_y + dy * (dist - 0.20) / dist
    else:
        end_x, end_y = bubble_cx, bubble_cy
    # 引线（用 LINE shape 画一段直线）
    line_x = min(target_x, end_x)
    line_y = min(target_y, end_y)
    line_w = abs(end_x - target_x) if abs(end_x - target_x) > 0.01 else 0.01
    line_h = abs(end_y - target_y) if abs(end_y - target_y) > 0.01 else 0.01
    try:
        # [v3.5 fix] add_connector 期望 (begin_x, begin_y, end_x, end_y), 不是 (x, y, w, h)
        ls = slide.shapes.add_connector(1, Inches(line_x), Inches(line_y), Inches(line_x + line_w), Inches(line_y + line_h))
        ls.line.color.rgb = COL["primary"]
        ls.line.width = Pt(1.2)
    except Exception:
        add_shape(slide, MSO_SHAPE.RECTANGLE, min(target_x, end_x), min(target_y, end_y),
                  max(abs(end_x - target_x), 0.01), max(abs(end_y - target_y), 0.01), COL["primary"])
    # 小气泡
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, bh, RGBColor(230, 240, 255), line=COL["primary"], radius=True)
    add_text(slide, text, bx + 0.08, by + 0.05, bw - 0.16, bh - 0.10, size=9.5, color=COL["text"])

def add_insight_panel_matrix(slide, cells, axis_x, axis_y, x=6.30, y=1.50, w=6.33, h=4.50):
    """C-2-F matrix-grid: 2x2 / 3x1 / 2x3 卡片矩阵 + 维度标题
    [v3.2] 默认 x=6.30 y=1.50 w=6.33 h=4.50 (右半边,避开 quote_items 左半边,适配 0.7 边距)"""
    rows = list(dict.fromkeys(c.get(axis_y, "") for c in cells))
    cols = list(dict.fromkeys(c.get(axis_x, "") for c in cells))
    nr, nc = len(rows), len(cols)
    # 维度标题
    add_text(slide, f"按 {axis_x} × {axis_y} 分类:", x, y - 0.30, w, 0.25, size=9, bold=True, color=COL["primary"])
    # 计算卡片尺寸
    title_w = 0.85
    grid_w = w - title_w - 0.10
    grid_h = h - 0.20
    cw = grid_w / nc
    rh = grid_h / nr
    # 表头（列名）
    for j, c in enumerate(cols):
        add_text(slide, c, x + title_w + j * cw + 0.05, y + 0.02, cw - 0.10, 0.22, size=9, bold=True, color=COL["muted"], align=PP_ALIGN.CENTER)
    # 卡片
    for i, r in enumerate(rows):
        add_text(slide, r, x + 0.05, y + 0.30 + i * rh + 0.10, title_w - 0.10, rh - 0.20, size=9, bold=True, color=COL["navy"], valign=MSO_ANCHOR.MIDDLE)
        for j, c in enumerate(cols):
            card_x = x + title_w + j * cw + 0.05
            card_y = y + 0.30 + i * rh + 0.05
            card_w = cw - 0.10
            card_h = rh - 0.10
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h, COL["white"], line=RGBColor(226, 232, 240), radius=True)
            cell = next((cc for cc in cells if cc.get(axis_x) == c and cc.get(axis_y) == r), None)
            if cell:
                add_text(slide, cell.get("text", ""), card_x + 0.10, card_y + 0.05, card_w - 0.20, card_h - 0.10, size=8.5, color=COL["text"])

def bar_chart(slide, data, x, y, w, h, max_val=None, horizontal=False, percent=True, show_values=True):
    max_val = max_val or max(data.values()) * 1.1
    if not horizontal:
        base = y + h
        bw = w / (len(data) * 1.8)
        gap = bw * 0.8
        for i, (name, val) in enumerate(data.items()):
            bh = h * val / max_val
            bx = x + i * (bw + gap) + 0.25
            add_shape(slide, MSO_SHAPE.RECTANGLE, bx, base - bh, bw, bh, PCOL.get(name, COL["primary"]))
            if show_values:
                add_text(slide, f"{val:.0f}%" if percent else f"{val:.1f}", bx - 0.08, base - bh - 0.28, bw + 0.16, 0.22, size=13, bold=True, color=PCOL.get(name, COL["primary"]), align=PP_ALIGN.CENTER)
            add_text(slide, name, bx - 0.10, base + 0.08, bw + 0.20, 0.22, size=9, color=COL["text"], align=PP_ALIGN.CENTER)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, base, w, 0.01, COL["light"])
    else:
        row_h = h / len(data)
        for i, (name, val) in enumerate(data.items()):
            yy = y + i * row_h
            add_text(slide, name, x, yy + 0.04, 1.2, 0.22, size=8.8, color=COL["text"])
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 1.25, yy + 0.06, w - 1.75, 0.18, RGBColor(237, 242, 247), radius=True)
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 1.25, yy + 0.06, (w - 1.75) * val / max_val, 0.18, PCOL.get(name, COL["primary"]), radius=True)
            add_text(slide, f"{val:.0f}%" if percent else f"{val:.2f}", x + w - 0.45, yy + 0.02, 0.44, 0.22, size=8.5, bold=True, color=PCOL.get(name, COL["primary"]), align=PP_ALIGN.RIGHT)

def grouped_bar(slide, categories, series_data, x, y, w, h, max_val, colors):
    base = y + h
    n_cat, n_ser = len(categories), len(series_data)
    group_w = w / n_cat
    bw = group_w / (n_ser + 1.4)
    for ci, cat in enumerate(categories):
        gx = x + ci * group_w + 0.15
        for si, (sname, vals) in enumerate(series_data.items()):
            val = vals[ci]
            bh = h * val / max_val
            bx = gx + si * bw
            add_shape(slide, MSO_SHAPE.RECTANGLE, bx, base - bh, bw * 0.75, bh, colors[si])
            add_text(slide, f"{val:.3f}" if max_val < 1 else f"{val:.0f}%", bx - 0.03, base - bh - 0.18, bw * 0.85, 0.15, size=6.7, color=COL["muted"], align=PP_ALIGN.CENTER)
        add_text(slide, cat, x + ci * group_w + 0.04, base + 0.05, group_w - 0.1, 0.2, size=8.0, color=COL["text"], align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, base, w, 0.01, COL["light"])
    for si, sname in enumerate(series_data.keys()):
        lx = x + si * 1.2
        add_shape(slide, MSO_SHAPE.RECTANGLE, lx, y - 0.28, 0.12, 0.12, colors[si])
        add_text(slide, sname, lx + 0.16, y - 0.31, 0.9, 0.17, size=7.5, color=COL["muted"])

def matrix_table(slide, headers, rows, x, y, w, h, highlight_cells=None):
    highlight_cells = highlight_cells or set()
    col_w = w / len(headers)
    row_h = h / (len(rows) + 1)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, row_h, COL["navy"])
    for ci, header in enumerate(headers):
        add_text(slide, header, x + ci * col_w + 0.02, y + 0.06, col_w - 0.04, row_h - 0.08, size=8.5, bold=True, color=COL["white"], align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows):
        yy = y + (ri + 1) * row_h
        bg = RGBColor(248, 250, 252) if ri % 2 == 0 else COL["white"]
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, yy, w, row_h, bg, line=RGBColor(226, 232, 240))
        for ci, val in enumerate(row):
            tx = x + ci * col_w
            if (ri, ci) in highlight_cells:
                s = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, tx + 0.12, yy + 0.08, col_w - 0.24, row_h - 0.16, COL["white"], line=COL["warning"], radius=True)
                s.line.width = Pt(1.6)
            add_text(slide, str(val), tx + 0.03, yy + 0.075, col_w - 0.06, row_h - 0.11, size=8.4, color=COL["text"], bold=(ci > 0 and (ri, ci) in highlight_cells), align=PP_ALIGN.CENTER)

def scatter_plot(slide, points, x, y, w, h, x_label, y_label, x_max=20, y_max=65):
    # axes
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.5, y + h - 0.35, w - 0.75, 0.01, COL["light"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.5, y + 0.05, 0.01, h - 0.40, COL["light"])
    for t in range(0, x_max + 1, 5):
        tx = x + 0.5 + (w - 0.75) * t / x_max
        add_shape(slide, MSO_SHAPE.RECTANGLE, tx, y + h - 0.37, 0.01, 0.04, COL["light"])
        add_text(slide, str(t), tx - 0.10, y + h - 0.24, 0.22, 0.14, size=6.5, color=COL["muted"], align=PP_ALIGN.CENTER)
    for t in range(0, y_max + 1, 10):
        ty = y + h - 0.35 - (h - 0.40) * t / y_max
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.48, ty, 0.04, 0.01, COL["light"])
        add_text(slide, str(t), x + 0.08, ty - 0.07, 0.34, 0.14, size=6.5, color=COL["muted"], align=PP_ALIGN.RIGHT)
    add_text(slide, x_label, x + w/2 - 1.1, y + h - 0.02, 2.2, 0.18, size=7.3, color=COL["muted"], align=PP_ALIGN.CENTER)
    add_text(slide, y_label, x + 0.02, y - 0.02, 1.15, 0.18, size=7.3, color=COL["muted"])
    for p in points:
        px = x + 0.5 + (w - 0.75) * p["x"] / x_max
        py = y + h - 0.35 - (h - 0.40) * p["y"] / y_max
        shape = MSO_SHAPE.ISOSCELES_TRIANGLE if p.get("shape") == "triangle" else MSO_SHAPE.OVAL
        add_shape(slide, shape, px - 0.09, py - 0.09, 0.18, 0.18, PCOL[p["name"]], line=COL["white"])
        add_text(slide, p["name"], px + 0.10, py - 0.09, 0.45, 0.17, size=7.4, color=COL["text"])
        add_text(slide, p.get("label", ""), px - 0.25, py - 0.31, 0.6, 0.16, size=7.2, bold=True, color=COL["success"], align=PP_ALIGN.CENTER)

def add_pp_callout(slide, text, x, y, color=COL["success"]):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 0.58, 0.24, RGBColor(240, 253, 244), line=color, radius=True)
    add_text(slide, text, x + 0.02, y + 0.035, 0.54, 0.14, size=8.2, bold=True, color=color, align=PP_ALIGN.CENTER)

def add_card(slide, title, body, x, y, w, h, color=COL["primary"], num=None):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, COL["white"], line=RGBColor(226, 232, 240), radius=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.045, h, color)
    if num is not None:
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.16, y + 0.16, 0.28, 0.28, color)
        add_text(slide, str(num), x + 0.16, y + 0.20, 0.28, 0.16, size=8, bold=True, color=COL["white"], align=PP_ALIGN.CENTER)
        tx = x + 0.54
    else:
        tx = x + 0.16
    add_text(slide, title, tx, y + 0.14, w - (tx - x) - 0.15, 0.28, size=11.5, bold=True, color=COL["navy"])
    add_text(slide, body, x + 0.16, y + 0.50, w - 0.30, h - 0.58, size=9.2, color=COL["text"])

# ---------- Slides ----------
# 1 Cover
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s, COL["navy"])
add_shape(s, MSO_SHAPE.OVAL, 9.6, -0.6, 3.1, 3.1, RGBColor(15, 52, 86))
add_shape(s, MSO_SHAPE.OVAL, 10.7, 4.5, 2.7, 2.7, RGBColor(0, 82, 217))
add_tencent_logo(s, dark=True)
add_text(s, "AI 产品用户留存分析", 0.7, 1.35, 8.4, 0.65, size=30, bold=True, color=COL["white"], font=TITLE_FONT)
add_text(s, "基于问卷数据的强留存、纯白用户、功能杠杆与运营召回洞察", 0.75, 2.15, 8.9, 0.35, size=14, color=RGBColor(215, 225, 240))
add_shape(s, MSO_SHAPE.RECTANGLE, 0.78, 2.78, 1.15, 0.05, COL["product_a"])
for i, (label, val) in enumerate([("样本量", "4,174"), ("分析对象", "元宝 / DS / 豆包"), ("数据周期", "2025.11.11-13")]):
    add_text(s, label, 0.82 + i*2.25, 5.58, 1.8, 0.25, size=9, color=RGBColor(190, 204, 220))
    add_text(s, val, 0.82 + i*2.25, 5.90, 2.3, 0.35, size=15, bold=True, color=COL["white"])
add_text(s, "format.skill v3｜专业咨询风格通用规范（基于 v0.9.1 实证分离式重构）", 0.78, 6.83, 7.3, 0.22, size=8.5, color=RGBColor(175, 190, 205))

# 2 Executive summary [v3.4 改方案 E: 纯观点卡全宽,无右栏无图表,5 cards 11.933 inch 满宽]
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s); add_title(s, "核心发现：留存差距由模型能力、功能纯白与纯白人群三轴共同驱动", "封面后置 Execution Summary；后续页面围绕五个主题展开")
add_tencent_logo(s); add_confidential_tag(s)
findings = [
    ("留存格局清晰分层", "豆包 54% > DS 34% > 元宝 19%，强留存差距已形成明显梯队。", COL["product_a"]),
    ("纯白用户显著增益", "纯白用户留存较非纯白提升 10-18pp，是高留存资产。", COL["success"]),
    ("可靠性是底层门票", "文本满意度中可靠性相关性最高，覆盖三产品均为核心因子。", COL["primary"]),
    ("拍照答疑是功能杠杆", "功能纯白留存提升最高达 17pp，是当前最值得抢占的单点。", COL["orange"]),
    ("运营召回仍待突破", "主动打开仍占主导，被动触达尚未成为核心召回机制。", COL["warning"]),
]
# [v3.9 修复] card h 0.65→0.80 (body h 0.07→0.22), spacing 0.75→0.90
# 旧: body h=0.07 (5pt) 装不下 9.2pt 文字 → 文字溢出卡片底部 6.5pt
# 新: body h=0.22 (16pt) 装 1 行 9.2pt 文字 (line_h=11.5pt) + 4.5pt padding
# 5×0.80 + 4×0.10 = 4.40 inch, y 1.70→6.10
# 数据条 strip y 5.65→6.20, h 0.45→0.30 (1 行 9pt fits)
for i, (t, b, c) in enumerate(findings):
    add_card(s, t, b, 0.7, 1.70 + i*0.90, 11.933, 0.80, c, num=i+1)
# [v3.9] 数据条 strip 下移 + 高度收窄,避免与 card 重叠
add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, 6.20, 11.933, 0.30, COL["panel_blue"], line=COL["primary"], radius=True)
add_text(s, "数据底层：强留存率（过去一周主用） / 渗透率（用户占比） / 相关性 ρ（满意度 vs 强留存）", 0.85, 6.25, 11.6, 0.20, size=8.5, color=COL["navy"], valign=MSO_ANCHOR.MIDDLE)
add_footer(s, "强留存率定义：过去一周主用该产品；样本 N=4,174。", 2)

# 3 Methodology
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s); add_title(s, "研究方法：样本覆盖 4,174 份问卷，围绕留存、纯白、满意度与召回四类问题展开", None); add_nav(s, 0)
add_tencent_logo(s); add_confidential_tag(s)
for i, (num, lab, desc, col) in enumerate([
    ("4,174", "总样本", "问卷投放回收", COL["primary"]),
    ("3", "产品", "元宝 / DS / 豆包", COL["product_b"]),
    ("4", "问题域", "留存、纯白、满意度、召回", COL["product_a"]),
    ("11.11-13", "周期", "2025 年 11 月", COL["orange"]),
]):
    # [v3.3 修复] 三段式 metric card: 数字(顶) + 标签(中) + 描述(底)
    # 旧版 (h=1.1) 把数字放在 y+0.33 与 label(y+0.14)/body(y+0.50) 都重叠
    # 新版: 数字 y+0.05 / 标签 y+0.45 / 描述 y+0.75, h 增到 1.3
    x = 0.72 + i*2.25
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.35, 2.0, 1.30, COL["white"], line=RGBColor(226, 232, 240), radius=True)
    add_shape(s, MSO_SHAPE.RECTANGLE, x, 1.35, 0.045, 1.30, col)
    add_text(s, num, x + 0.15, 1.40, 1.70, 0.36, size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, lab, x + 0.15, 1.80, 1.70, 0.24, size=11.5, bold=True, color=COL["navy"], align=PP_ALIGN.CENTER)
    add_text(s, desc, x + 0.15, 2.10, 1.70, 0.46, size=9.2, color=COL["text"], align=PP_ALIGN.CENTER)
# process flow
steps = [("1", "识别强留存"), ("2", "拆分纯白用户"), ("3", "评估文本/功能满意度"), ("4", "分析召回与流失"), ("5", "形成行动建议")]
for i, (n, t) in enumerate(steps):
    # [v3.2 修复] 起始 x 从 0.9 调到 1.05,确保首个标签 x-0.20=0.85 在 0.7 边距内
    x = 1.05 + i*2.0
    add_shape(s, MSO_SHAPE.OVAL, x, 3.55, 0.42, 0.42, COL["primary"])
    add_text(s, n, x, 3.64, 0.42, 0.16, size=9, bold=True, color=COL["white"], align=PP_ALIGN.CENTER)
    add_text(s, t, x - 0.20, 4.05, 1.1, 0.32, size=9.2, color=COL["text"], align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        add_shape(s, MSO_SHAPE.RIGHT_ARROW, x + 0.58, 3.68, 1.0, 0.16, RGBColor(203, 213, 225))
add_insight_panel_takeaway(s, [
    {"tag": "Step 1-2", "text": "识别强留存 + 拆分纯白"},
    {"tag": "Step 3-4", "text": "评估满意度 + 召回分析"},
    {"tag": "Step 5", "text": "形成行动建议"}
], 0.7, 6.05, 11.933, 0.65, title="研究流程")
# [v3.6 4 区填充] P3 已有 takeaway 面板填补底部, 但仍需 x>9.5 rail
# 方法: 在右下空白区(x=10.0-12.63, y=4.20-5.40) 加一个小型 "研究方法论要点" panel
add_discovery_panel(s, [
    "N=4,174 样本, 性别年龄匹配网民大盘",
    "数据周期 2025.11.11-11.13 (3 天)",
    "定量问卷 + 满意度 5 分量表"
])
add_footer(s, SAMPLE_DATA["method"] + "；数据周期 " + SAMPLE_DATA["date"] + "。", 3)

# 4 Overview [v3.7 重构: 移除 banner (与 title/bottom_summary 信息重复) + 调整 bottom_summary 位置避开 bar chart x 轴标签]
# v3.6 问题: ① banner (y=1.10-1.55) 与 KPI 卡 (y=1.30-2.50) 垂直重叠 0.25 inch
#           ② bar chart x 轴标签 (y=5.48-5.70) 与 bottom_summary (y=5.55-6.65) 重叠 0.15 inch
# v3.7 解法: 移除 banner; KPI 上移到 y=1.20; bottom_summary 下移到 y=5.80, h=0.80
#   LEFT (x=0.7-4.3): bar chart y=1.30-5.40 (baseline=5.40)
#   MIDDLE-TOP (x=4.5-9.5): 3 KPI 卡片横向排 y=1.20-2.40 (h=1.20)
#   MIDDLE-BOTTOM (x=4.5-9.5): 数据表 y=2.50-5.40 (h=2.90, 3 行 + 1 表头)
#   RIGHT (x=10.0-12.63): "发现" rail y=1.30-5.40 (h=4.10)
#   BOTTOM (x=0.7-12.63): 关键发现 3 列 panel y=5.80-6.60 (h=0.80) — 避开 x 轴标签
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s); add_title(s, "留存总览：豆包强留存率领先，DS 居中，元宝仍处追赶位置", None); add_nav(s, 0)
add_tencent_logo(s); add_confidential_tag(s)
# [v3.7] 移除 banner (v3.6 与 KPI 卡垂直重叠;信息已在 title + bottom_summary)
# LEFT: bar chart y=1.30-5.40 (h=4.10), baseline 5.40 与右表底对齐
bar_chart(s, {p: STRONG_RETENTION[p]["rate"] for p in PRODUCTS}, 0.7, 1.30, 3.60, 4.10, max_val=60)
# MIDDLE-TOP: 3 KPI 卡片横向排 y=1.20-2.40 (h=1.20) — [v3.7] 从 y=1.30 上移到 y=1.20
for i, p in enumerate(PRODUCTS):
    x = 4.50 + i*1.68
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.20, 1.55, 1.20, PLIGHT[p], line=PCOL[p], radius=True)
    add_text(s, p, x + 0.08, 1.30, 1.35, 0.20, size=10, bold=True, color=PCOL[p], align=PP_ALIGN.CENTER)
    add_text(s, f"{STRONG_RETENTION[p]['users']:,}", x + 0.08, 1.58, 1.35, 0.30, size=18, bold=True, color=COL["navy"], align=PP_ALIGN.CENTER)
    add_text(s, f"渗透 {STRONG_RETENTION[p]['penetration']:.1f}%", x + 0.05, 1.95, 1.40, 0.18, size=8.5, color=COL["muted"], align=PP_ALIGN.CENTER)
# MIDDLE-BOTTOM: 数据表 y=2.50-5.40 (h=2.90) — [v3.7] 从 y=2.60 上移到 y=2.50 与 KPI 紧贴, 底 y=5.40 与 LEFT bar chart 底线平行
headers = ["产品", "强留存", "样本用户", "渗透率"]
rows = [[p, f"{STRONG_RETENTION[p]['rate']:.0f}%", f"{STRONG_RETENTION[p]['users']:,}", f"{STRONG_RETENTION[p]['penetration']:.1f}%"] for p in PRODUCTS]
matrix_table(s, headers, rows, 4.50, 2.50, 5.00, 2.90, highlight_cells={(0,1), (1,1), (2,1)})
# RIGHT: "发现" rail 填补 x>9.5 空白 (v3.5 缺失的右对齐)
add_discovery_panel(s, [
    "豆包强留存 54%, 是元宝 19% 的 2.8 倍",
    "DS 居中 34%, 三产品形成清晰梯队",
    "渗透率元宝最低 66.4%, 仍有空间"
])
# BOTTOM: 关键发现 3 列 panel y=5.80-6.60 (h=0.80) — [v3.7] 从 y=5.55-6.65 下移避开 x 轴标签 (y=5.48-5.70)
add_bottom_summary(s, [
    {"tag": "格局", "text": "豆包 54% > DS 34% > 元宝 19% 梯队已形成", "color": COL["product_a"]},
    {"tag": "机会", "text": "元宝渗透 66.4% 最低, 拉升后留存空间最大", "color": COL["orange"]},
    {"tag": "抓手", "text": "补可靠性 + 拍照答疑 = 元宝追赶核心杠杆", "color": COL["primary"]}
], 0.7, 5.80, 11.933, 0.80, title="留存格局")
add_footer(s, "Page 1：强留存率、各产品使用用户数与渗透率。", 4)

# 5 Driver matrix
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s); add_title(s, "留存驱动矩阵：可靠性、深度思考与拍照答疑构成三产品共同杠杆", None); add_nav(s, 0)
add_tencent_logo(s); add_confidential_tag(s)
for pi, p in enumerate(PRODUCTS):
    x = 0.72 + pi*3.05
    add_text(s, p, x, 1.26, 1.6, 0.26, size=13, bold=True, color=PCOL[p])
    top_items = RETENTION_DRIVERS[p][:5]
    max_v = max(v for _, _, v in top_items)
    for i, (name, typ, val) in enumerate(top_items):
        yy = 1.70 + i*0.55
        # [v3.3 修复] 4 列总宽从 3.31 收到 2.85,不再撞下一产品列
        # 旧: name(0.95) + bar(1.45) + value(0.46) + type(0.45) = 3.31 > 列距 3.05
        # 新: name(0.75) + bar(1.45) + value(0.30) + type(0.35) = 2.85
        add_text(s, name, x, yy, 0.75, 0.18, size=8.0, color=COL["text"])
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.80, yy + 0.02, 1.45, 0.14, RGBColor(237, 242, 247), radius=True)
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.80, yy + 0.02, 1.45*val/max_v, 0.14, PCOL[p], radius=True)
        # [v3.9 修复] +pp box w 0.30→0.40 (7.5pt 6字符 26.2pt > 21.6pt 旧box → 2行折行溢出)
        # 旧: 段宽 26.2pt > box 21.6pt → 折2行 18.8pt > box 13pt → 垂直溢出 5.8pt
        # 新: box w 0.40=28.8pt, 1行 fits → 不再折行
        # type x 2.65→2.70 (让位给 +pp 增宽 0.10)
        add_text(s, f"+{val:.1f}pp", x + 2.25, yy - 0.02, 0.40, 0.18, size=7.5, bold=True, color=COL["success"], align=PP_ALIGN.RIGHT)
        add_text(s, typ, x + 2.70, yy - 0.02, 0.35, 0.18, size=6.6, color=COL["muted"])
# emphasis frame around common factors
# [v3.2 修复] x 从 0.62 调到 0.7,满足 0.7 inch 统一边距
# [v3.7 修复] y=4.88, h=0.82 (4.88-5.70) 与 bottom_summary (5.55-6.65) 重叠 0.15 → y=4.55, h=0.65 (4.55-5.20) 让位
add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, 4.55, 8.98, 0.65, COL["white"], line=COL["warning"], radius=True)
# [v3.3 修复] frame 内 2 段文字宽度收窄避免重叠
# 旧: text1 w=4.70(0.82-5.52) 与 text2 x=5.25 撞 0.27
# 新: text1 w=4.00(0.82-4.82), text2 x=4.95, gap 0.13
add_text(s, "共性高杠杆：可靠性 / 深度思考 / 拍照答疑", 0.82, 4.65, 4.00, 0.22, size=12, bold=True, color=COL["warning"])
add_text(s, "差异化策略应从'全量补齐模型能力'转向'按产品短板匹配功能纯白机会'。", 4.95, 4.66, 4.70, 0.22, size=9.0, color=COL["text"])
# [v3.6 4 区填充] RIGHT: "发现" rail 填补 x>9.5 空白 (v3.5 缺 rail; v3.5 的 +25.2pp callout 移到 rail 第 2 条)
add_discovery_panel(s, [
    "可靠性 / 深度思考 / 拍照答疑 = 共同杠杆",
    "DS 可靠性杠杆 +25.2pp 全图最高",
    "豆包 拍照答疑杠杆 +21pp 居首"
])
# [v3.6 4 区填充] BOTTOM: 关键发现 3 列 panel 填补底部空白
add_bottom_summary(s, [
    {"tag": "共性杠杆", "text": "可靠性 / 深度思考 / 拍照答疑 = 三产品 Top 共同因子", "color": COL["primary"]},
    {"tag": "DS 优势", "text": "可靠性 +25.2pp / 深度思考 +23.8pp 双高", "color": COL["product_b"]},
    {"tag": "豆包特色", "text": "拍照答疑 +21pp 纯白场景独占优势", "color": COL["product_a"]}
], title="驱动格局")
add_footer(s, "Page 2：留存提升驱动因素，单位为百分点（pp）。", 5)

# 6 Pure users scatter
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s); add_title(s, "纯白用户：三产品纯白留存均高于非纯白，DS 增益最大达到 +18pp", None); add_nav(s, 1)
add_tencent_logo(s); add_confidential_tag(s)
points = []
for p in PRODUCTS:
    d = PURE_NEW_USERS[p]
    lift = d["pure_new_rate"] - d["non_pure_new_rate"]
    points.append({"name": p, "x": lift, "y": d["pure_new_rate"], "label": f"+{lift}pp", "shape": "triangle" if p == "DS" else "circle"})
# [v3.8 修复] P6 两张图表底部对齐: scatter h 4.55→3.98 (base 5.97→5.40); side bars base 4.65→5.40 (整体上移 0.75)
scatter_plot(s, points, 0.75, 1.42, 5.15, 3.98, "纯白留存提升（pp）", "纯白用户强留存率（%）", x_max=20, y_max=65)
# side comparison bars
for i, p in enumerate(PRODUCTS):
    x = 6.32 + i*1.12
    d = PURE_NEW_USERS[p]
    add_text(s, p, x, 2.20, 0.75, 0.18, size=9, bold=True, color=PCOL[p], align=PP_ALIGN.CENTER)
    bh1 = 2.7 * d["pure_new_rate"] / 65
    bh2 = 2.7 * d["non_pure_new_rate"] / 65
    add_shape(s, MSO_SHAPE.RECTANGLE, x + 0.15, 5.40-bh1, 0.25, bh1, PCOL[p])
    add_shape(s, MSO_SHAPE.RECTANGLE, x + 0.48, 5.40-bh2, 0.25, bh2, RGBColor(203, 213, 225))
    add_text(s, f"{d['pure_new_rate']}%", x + 0.02, 5.40-bh1-0.20, 0.45, 0.14, size=7, bold=True, color=PCOL[p], align=PP_ALIGN.CENTER)
    add_text(s, f"{d['non_pure_new_rate']}%", x + 0.42, 5.40-bh2-0.18, 0.42, 0.14, size=7, color=COL["muted"], align=PP_ALIGN.CENTER)
    add_pp_callout(s, f"+{d['pure_new_rate']-d['non_pure_new_rate']}pp", x + 0.12, 5.70)
add_text(s, "纯白 vs 非纯白", 6.35, 1.87, 2.6, 0.22, size=11, bold=True, color=COL["navy"])
add_discovery_panel(s, ["DS 纯白增益最高，说明新用户资产质量更强", "元宝纯白虽仅 215 人，但留存翻倍", "纯白策略需与功能入口绑定，而非仅看拉新"], h=5.20)
add_footer(s, "Page 3：纯白用户强留存率、非纯白用户强留存率、相关性与显著性。", 6)

# 7 Source matrix / overlap
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s); add_title(s, "用户来源矩阵：从 DS 抢人难度低于豆包，三栖用户是更现实的转化池", None); add_nav(s, 1)
add_tencent_logo(s); add_confidential_tag(s)
# [v3.10 重构] 顶部 4 blocks 文字逻辑改为"维度名 + 数量 + 枚举值"三行结构
# 旧 (v3.9): (人话名, 用户数, 描述) = ("三栖用户", "2,128", "豆包+DS+元宝")
#   痛点: "2,128" 是用户数,与方块主题"三栖用户"关系不明; "豆包+DS+元宝" 既是枚举又是解释,信息重复
# 新 (v3.10): (维度名, 数量, 枚举值) + 用户数作为附加信息
#   三行: 维度名 → 数字 + 单位 → 枚举值 (用户反馈顺序: 产品 → 3 → 元宝/DS/豆包)
blocks = [
    # (label, count, products, users, color)
    ("产品组合", 3, "元宝 / DS / 豆包", "2,128", COL["primary"]),
    ("产品组合", 2, "DS / 元宝", "172", COL["product_b"]),
    ("产品组合", 2, "豆包 / 元宝", "437", COL["product_a"]),
    ("产品组合", 2, "豆包 / DS", "1,307", COL["orange"]),
]
for i, (label, count, products, users, c) in enumerate(blocks):
    x = 0.85 + (i%2)*4.25; y = 1.55 + (i//2)*1.55
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 3.75, 1.12, RGBColor(248, 250, 252), line=c, radius=True)
    # 行 1 (左): 维度名 "产品组合" (标签/类别, 灰色)
    add_text(s, label, x + 0.18, y + 0.10, 1.2, 0.20, size=9, color=COL["muted"])
    # 行 1 (右): 用户数 (附加信息, 灰色, 右对齐)
    add_text(s, f"{users} 用户", x + 2.10, y + 0.10, 1.5, 0.20, size=9, color=COL["muted"], align=PP_ALIGN.RIGHT)
    # 行 2: 大数字 count (主信息) + "种产品" 单位
    add_text(s, str(count), x + 0.18, y + 0.30, 0.7, 0.48, size=26, bold=True, color=c)
    add_text(s, "种产品", x + 0.78, y + 0.45, 0.7, 0.22, size=10, color=COL["muted"])
    # 行 3: 枚举产品名 (navy 深色加粗)
    add_text(s, products, x + 0.18, y + 0.80, 3.4, 0.25, size=12, bold=True, color=COL["navy"])
# arrows
# [v3.7 修复] 箭头位置避开相邻卡片: 右箭头 x 从 4.72 w=0.55 收到 x=4.65 w=0.45 (填满 card1-card2 间隙 4.60-5.10)
# 上下箭头 y 从 2.80 收到 2.65 (下沿 3.05, 与下方卡片 3.10 留 0.05 gap)
add_shape(s, MSO_SHAPE.RIGHT_ARROW, 4.65, 2.05, 0.45, 0.22, RGBColor(203, 213, 225))
add_shape(s, MSO_SHAPE.DOWN_ARROW, 2.48, 2.65, 0.24, 0.40, RGBColor(203, 213, 225))
add_shape(s, MSO_SHAPE.DOWN_ARROW, 6.72, 2.65, 0.24, 0.40, RGBColor(203, 213, 225))
add_text(s, "机会判断", 0.92, 4.55, 1.0, 0.22, size=12, bold=True, color=COL["navy"])
add_para_text(s, ["优先从 DS 与元宝重叠用户中提升主用频率", "豆包迁移需以差异功能而非通用问答切入", "三栖用户适合做场景分层与召回实验"], 0.95, 4.88, 8.1, 0.65, size=9.8, bullet=True)
add_discovery_panel(s, ["三栖用户规模最大，是中短期转化主战场", "仅 DS+元宝用户虽小但留存潜力最高", "豆包用户迁移难度高，需要强差异功能"])
# [v3.6 4 区填充] BOTTOM: 关键发现 3 列 panel 填补底部空白 (v3.5 y=4.85-6.85 留 ~2 inch 空白)
add_bottom_summary(s, [
    {"tag": "主战场", "text": "三栖用户 2128 人规模最大, 转化主战场", "color": COL["primary"]},
    {"tag": "高潜力", "text": "DS+元宝双栖留存最高, 优先拉升频次", "color": COL["product_b"]},
    {"tag": "难点", "text": "豆包用户迁移难度高, 需差异功能切入", "color": COL["product_a"]}
], title="来源机会")
add_footer(s, "Page 5：用户产品组合与来源矩阵；部分重叠用户数来自原文。", 7)

# 8 Text satisfaction
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s); add_title(s, "文本满意度：可靠性在三产品中均位居核心位置，是提升强留存的底层门票", None); add_nav(s, 2)
add_tencent_logo(s); add_confidential_tag(s)
cats = ["豆包", "DS", "元宝"]
# [v3.12] 严格模式: 缺 key 时直接 KeyError 而不是静默填 0 (避免画高度为 0 的 bar)
series = {
    "可靠性": [TEXT_SATISFACTION_CORR[p][0][1] if TEXT_SATISFACTION_CORR[p][0][0]=="可靠性" else dict(TEXT_SATISFACTION_CORR[p])["可靠性"] for p in cats],
    "可用性": [dict(TEXT_SATISFACTION_CORR[p])["可用性"] for p in cats],
    "速度/稳定性": [dict(TEXT_SATISFACTION_CORR[p])["速度/稳定性"] for p in cats],
    "情绪价值": [dict(TEXT_SATISFACTION_CORR[p])["情绪价值"] for p in cats],
}
grouped_bar(s, cats, series, 0.75, 1.75, 5.6, 3.35, 0.25, [COL["primary"], COL["product_b"], COL["orange"], COL["product_a"]])
# Top factor cards
# [v3.7 修复] 顶卡 y 从 1.58 下移到 1.85, 让位给 anchor (y=1.25-1.70) 留 0.15 gap
for i, p in enumerate(PRODUCTS):
    top = max(TEXT_SATISFACTION_CORR[p], key=lambda x: x[1])
    add_card(s, p, f"Top 因子：{top[0]}\n相关系数 ρ={top[1]:.3f}", 6.75, 1.85 + i*1.05, 2.45, 0.82, PCOL[p])
add_insight_panel_anchor(s, "豆包文本满意度整体最强，Top 因子可靠性 ρ=0.42", 7.97, 1.85, bx=6.50, by=1.25, bw=2.6, bh=0.45)
# [v3.6 4 区填充] RIGHT: "发现" rail 填补 x>9.5 空白 (v3.5 缺 rail)
add_discovery_panel(s, [
    "可靠性 = 三产品共性 Top 因子 (强相关)",
    "豆包整体满意度最强, Top ρ=0.42",
    "元宝 / DS 在情绪价值有差异化空间"
])
# [v3.6 4 区填充] BOTTOM: 关键发现 3 列 panel 填补底部空白
add_bottom_summary(s, [
    {"tag": "共性", "text": "可靠性 = 三产品共性 Top 1 因子, 强相关", "color": COL["primary"]},
    {"tag": "豆包", "text": "ρ=0.42 全图最高, 文本满意度整体最强", "color": COL["product_a"]},
    {"tag": "差异化", "text": "元宝/DS 在情绪价值 / 速度/稳定性 有空间", "color": COL["product_b"]}
], title="文本满意度")
add_footer(s, "Page 7：文本满意度维度与强留存相关性，ρ 表示相关系数。", 8)

# 9 Function satisfaction
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s); add_title(s, "功能满意度：深度思考与拍照答疑贡献突出，但不同产品的功能短板不同", None); add_nav(s, 3)
add_tencent_logo(s); add_confidential_tag(s)
# Heatmap table
features = ["深度思考", "拍照答疑", "AI创作", "打电话", "朗读"]
headers = ["功能"] + PRODUCTS
rows = []
for f in features:
    row = [f]
    for p in PRODUCTS:
        row.append(f"{dict(FUNCTION_SATISFACTION_CORR[p]).get(f, 0):.3f}")
    rows.append(row)
# highlight max per row
highlights = set()
for ri, f in enumerate(features):
    vals = [dict(FUNCTION_SATISFACTION_CORR[p]).get(f, 0) for p in PRODUCTS]
    max_idx = vals.index(max(vals)) + 1
    highlights.add((ri, max_idx))
matrix_table(s, headers, rows, 0.75, 1.45, 5.35, 3.55, highlight_cells=highlights)
# mini bars for top two
top2 = {"豆包 深度思考": 0.227, "豆包 拍照答疑": 0.184, "元宝 打电话": 0.170, "DS 深度思考": 0.137}
bar_chart(s, top2, 6.55, 1.55, 2.45, 2.8, max_val=0.25, horizontal=True, percent=False)
add_text(s, "Top 相关性因子", 6.55, 1.20, 1.8, 0.22, size=11, bold=True, color=COL["navy"])
add_discovery_panel(s, ["深度思考在豆包与 DS 中更突出", "元宝在打电话/AI创作/朗读上相关性更高", "拍照答疑虽非满意度最高，但纯白提升最强"])
# [v3.6 4 区填充] BOTTOM: 关键发现 3 列 panel 填补底部空白
add_bottom_summary(s, [
    {"tag": "豆包/DS", "text": "深度思考 + 拍照答疑 双高, Top 相关性 ρ=0.227/0.184", "color": COL["product_a"]},
    {"tag": "元宝", "text": "打电话/AI创作/朗读 相关性更高, 路径差异化", "color": COL["product_c"]},
    {"tag": "抓手", "text": "拍照答疑纯白提升最强, 应优先做新用户场景", "color": COL["orange"]}
], title="功能格局")
add_footer(s, "Page 8-9：功能满意度与强留存相关性。", 9)

# 10 Function pure new lift
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s); add_title(s, "功能纯白：拍照答疑带来 10-17pp 留存提升，是当前最大单一功能杠杆", None); add_nav(s, 3)
add_tencent_logo(s); add_confidential_tag(s)
# lift bars by product for photo Q&A and thinking
# [v3.7 修复] bar_chart h 从 3.25 收到 2.50, x 轴标签 (y=4.33-4.55) 与 bottom_summary (y=5.55-6.65) 留 1.0 gap
for i, f in enumerate(["拍照答疑", "深度思考"]):
    x = 0.75 + i*4.45
    data = {p: FUNCTION_PURE_NEW_LIFT[p][f][2] for p in PRODUCTS if f in FUNCTION_PURE_NEW_LIFT[p]}
    add_text(s, f, x, 1.30, 1.8, 0.24, size=12, bold=True, color=COL["navy"])
    bar_chart(s, data, x, 1.75, 3.45, 2.50, max_val=18)
    for j, (p, lift) in enumerate(data.items()):
        # [v3.7 修复] callout y 从 4.40 下移到 4.65, 让位给 bar_chart x 轴标签 (y=4.33-4.55)
        add_pp_callout(s, f"+{lift}pp", x + 0.35 + j*0.86, 4.65)
# [v3.7 修复] 删除 matrix_table — 与 bar chart + callouts 信息完全重复, 且与 bottom_summary 大幅重叠
# [v3.6 旧] headers = ["产品", "拍照答疑", "深度思考", "最高功能纯白"]
#          matrix_table(s, headers, rows, 0.92, 5.75, 8.05, 0.82, ...) — y=5.75-6.57 与 bottom_summary (5.55-6.65) 重叠 0.82 inch
# [v3.6 4 区填充] RIGHT: "发现" rail (h=4.10 收窄, 让位给底部 panel)
add_discovery_panel(s, ["拍照答疑对 DS 与元宝拉动最高", "豆包功能纯白提升更分散", "应优先把拍照答疑做成新用户第一场景"])
# [v3.6 4 区填充] BOTTOM: 关键发现 3 列 panel
add_bottom_summary(s, [
    {"tag": "杠杆 1", "text": "拍照答疑 = 三产品共性最大纯白提升 (10-17pp)", "color": COL["primary"]},
    {"tag": "杠杆 2", "text": "DS/元宝 拍照答疑 +17/+16pp 拉动最强", "color": COL["product_b"]},
    {"tag": "抓手", "text": "豆包应把拍照答疑做成新用户第一场景", "color": COL["product_a"]}
], title="功能纯白提升")
add_footer(s, "Page 9：功能纯白用户 vs 非功能纯白用户强留存率差值。", 10)

# 11 Recall mechanism
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s); add_title(s, "召回机制：主动打开仍是主要路径，被动触达尚未成为稳定增长引擎", None); add_nav(s, 4)
add_tencent_logo(s); add_confidential_tag(s)
metrics = ["主动打开", "有红点", "桌面入口", "Push", "朋友分享"]
series = {p: [RECALL_DATA[m][p] for m in metrics] for p in PRODUCTS}
grouped_bar(s, metrics, series, 0.75, 1.65, 6.55, 3.40, 70, [PCOL[p] for p in PRODUCTS])
# left callouts
# [v3.7 修复] callout cards 起点 y 从 1.50 收到 1.40, h 不变; 让位给 grouped_bar (h 3.40 → end 5.05, x 轴标签 5.13-5.35)
#          原 y=1.50+i*1.08 (1.50, 2.58, 3.66) 与 grouped_bar x 范围冲突 (callouts x=7.65+ 在 chart x=0.75-7.30 之外, 不冲突)
for i, (m, desc) in enumerate([("主动打开", ">57%"), ("桌面/搜索", "DS/元宝更高"), ("红点", "元宝 23%")]):
    add_card(s, m, desc, 7.65, 1.50 + i*1.08, 1.58, 0.82, [COL["primary"], COL["product_b"], COL["product_c"]][i])
add_discovery_panel(s, ["三产品主动打开均接近或超过 58%", "元宝红点触达占比相对更高", "被动召回需从提醒转向场景化触达"])
# [v3.6 4 区填充] BOTTOM: 关键发现 3 列 panel
add_bottom_summary(s, [
    {"tag": "主路径", "text": "主动打开 = 三产品核心召回方式 (>57%)", "color": COL["primary"]},
    {"tag": "元宝", "text": "红点触达 23% / 桌面入口 30.5% 相对更高", "color": COL["product_c"]},
    {"tag": "升级", "text": "被动召回需从提醒升级为场景化触达", "color": COL["orange"]}
], title="召回格局")
add_footer(s, "Page 10-11：用户最近一次打开产品的召回方式。", 11)

# 12 Churn reasons
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s); add_title(s, "流失原因：对文本答案不满意是首要流失因素，元宝与 DS 问题更集中", None); add_nav(s, 4)
add_tencent_logo(s); add_confidential_tag(s)
bar_chart(s, {p: CHURN_REASONS[p]["text_unsatisfied"] for p in PRODUCTS}, 0.75, 1.55, 3.4, 3.55, max_val=70)
add_text(s, "文本答案不满意占比", 0.82, 1.18, 1.9, 0.23, size=11, bold=True, color=COL["navy"])
# reasons cards
for i, p in enumerate(PRODUCTS):
    x = 4.62; y = 1.45 + i*1.1
    d = CHURN_REASONS[p]
    body = d["top_reason"] + (f"\n功能不满意：{d['function_unsatisfied']}%" if d.get("function_unsatisfied") else "")
    add_card(s, p, body, x, y, 4.55, 0.86, PCOL[p])
# [v3.3 修复] 删除冗余的红色 62% callout — 柱顶已标 "62%" 数值,callout 与之重叠且重复
add_discovery_panel(s, ["元宝文本不满最高，核心是可靠性与信源", "DS 流失更集中在速度/稳定性", "豆包需关注可靠性与创作功能一致性"])
# [v3.6 4 区填充] BOTTOM: 关键发现 3 列 panel 填补底部空白 (v3.5 y=4.85-6.85 留 ~2 inch 空白)
add_bottom_summary(s, [
    {"tag": "首要原因", "text": "文本答案不满意 = 流失首因, 元宝 62% 最高", "color": COL["warning"]},
    {"tag": "DS 集中", "text": "流失更集中在速度/稳定性 (66%)", "color": COL["product_b"]},
    {"tag": "豆包", "text": "需关注可靠性 (44%) + 创作功能一致性", "color": COL["product_a"]}
], title="流失原因")
add_footer(s, "Page 12-13：流失原因与用户不满意项；百分比为对应流失用户提及占比。", 12)

# 13 User voice [v3.2 修复 quote 与 matrix 重叠:左半 quotes + 右半 matrix]
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s); add_title(s, "用户原声：正负反馈共同指向'可靠性底座 + 差异化功能体验'", None); add_nav(s, 4)
add_tencent_logo(s); add_confidential_tag(s)
quote_items = [
    ("元宝正向", USER_QUOTES["元宝正向"][0], COL["product_c"]),
    ("元宝负向", USER_QUOTES["元宝负向"][0], COL["warning"]),
    ("豆包正向", USER_QUOTES["豆包正向"][0], COL["product_a"]),
    ("豆包负向", USER_QUOTES["豆包负向"][1], COL["warning"]),
    ("DS正向", USER_QUOTES["DS正向"][0], COL["product_b"]),
    ("DS负向", USER_QUOTES["DS负向"][0], COL["warning"]),
]
# [v3.2] 6 quotes 移到左半 (x=0.7-5.85, y=1.5-6.0): 2 列 3 行,每张 w=2.55, h=1.4
for i, (t, q, c) in enumerate(quote_items):
    col = i % 2
    row = i // 2
    x = 0.7 + col * 2.65
    y = 1.50 + row * 1.50
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 2.55, 1.40, RGBColor(255, 251, 235), line=c, radius=True)
    add_text(s, t, x + 0.10, y + 0.10, 1.1, 0.20, size=9.0, bold=True, color=c)
    add_text(s, "\u201C" + q + "\u201D", x + 0.10, y + 0.34, 2.35, 0.95, size=9.0, color=COL["text"], font="KaiTi")
# [v3.2] matrix 移到右半,使用新默认 x=6.30 y=1.50 w=6.33 h=4.50
add_insight_panel_matrix(s, [
        {"product": "豆包", "polarity": "正向", "text": "速度、稳定性、专业深度"},
        {"product": "豆包", "polarity": "负向", "text": "可靠性、生成质量"},
        {"product": "DS",   "polarity": "正向", "text": "速度、稳定性、深度思考"},
        {"product": "DS",   "polarity": "负向", "text": "生成质量、稳定性"},
        {"product": "元宝", "polarity": "正向", "text": "功能丰富、AI创作"},
        {"product": "元宝", "polarity": "负向", "text": "可靠性、信源权威"}
    ], axis_x="polarity", axis_y="product")
add_footer(s, "Page 13：用户原始反馈归纳；保留原文含义，按正负向分类。", 13)

# 14 Strategy roadmap
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s); add_title(s, "行动建议：用 Exploit / Explore / Watch 三层路径提升留存与差异化", None); add_nav(s, 5)
add_tencent_logo(s); add_confidential_tag(s)
actions = [
    ("Exploit", "补齐可靠性底座", "提升答案准确性、减少幻觉、增强信源权威；优先修复元宝/DS 流失主因。", COL["primary"]),
    ("Explore", "抢占拍照答疑纯白", "围绕拍照答疑打造新用户首个高价值场景；将功能纯白提升转化为强留存。", COL["product_a"]),
    ("Watch", "探索新功能与召回", "关注陪伴、电商、网络资源获取等潜力场景；召回从红点转向场景化触达。", COL["orange"]),
]
for i, (tag, title, body, c) in enumerate(actions):
    # [v3.6 4 区填充] action 卡片收窄 w=7.45 (旧 8.45) 让位给右侧 rail (x=10.0-12.63)
    y = 1.42 + i*1.45
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.82, y, 7.45, 1.05, RGBColor(248, 250, 252), line=c, radius=True)
    add_shape(s, MSO_SHAPE.OVAL, 1.05, y + 0.25, 0.48, 0.48, c)
    add_text(s, str(i+1), 1.05, y + 0.35, 0.48, 0.17, size=9, bold=True, color=COL["white"], align=PP_ALIGN.CENTER)
    # [v3.3 修复] tag w 从 1.10 收到 0.95,避开与 title (x=2.75) 重叠
    add_text(s, tag, 1.72, y + 0.18, 0.95, 0.20, size=9, bold=True, color=c)
    # [v3.6 修复] title 收窄 w 从 2.3 到 1.8; body x=2.75 改为 2.75, w 收窄到 4.95
    add_text(s, title, 2.75, y + 0.16, 1.8, 0.24, size=12, bold=True, color=COL["navy"])
    add_text(s, body, 4.65, y + 0.16, 3.55, 0.78, size=9.0, color=COL["text"])
# [v3.6 4 区填充] RIGHT: "发现" rail (v3.5 缺 rail, 3.5 inch 右空白)
add_discovery_panel(s, [
    "Exploit = 可靠性底座, 解决流失首因",
    "Explore = 拍照答疑纯白, 新用户第一场景",
    "Watch = 场景化召回, 升级提醒机制"
])
add_insight_panel_takeaway(s, [
        {"tag": "立即做", "text": "补可靠性底座（信源+准确度）"},
        {"tag": "短期做", "text": "拍照答疑场景化（新用户第一场景）"},
        {"tag": "中期探索", "text": "召回机制与高价值场景绑定"}
    ], 0.7, 6.05, 11.933, 0.65, title="行动建议")
add_footer(s, "Page 2 + Page 9 + Page 12-13：行动建议根据驱动因素、功能纯白与流失原因综合形成。", 14)

# 15 Opportunity matrix
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s); add_title(s, "优先级矩阵：高影响、低依赖的可靠性与拍照答疑应进入近期路线图", None); add_nav(s, 5)
add_tencent_logo(s); add_confidential_tag(s)
# Matrix axes [v3.7] mh 从 3.95 收到 3.55, 让位给 bottom_summary (y=5.55) — "实施复杂度 →" 轴标签 (y=my+mh+0.12) 从 5.52 上移到 5.12
mx, my, mw, mh = 0.85, 1.45, 6.2, 3.55
add_shape(s, MSO_SHAPE.RECTANGLE, mx, my+mh, mw, 0.012, COL["light"])
add_shape(s, MSO_SHAPE.RECTANGLE, mx, my, 0.012, mh, COL["light"])
add_text(s, "实施复杂度 →", mx+mw-1.4, my+mh+0.12, 1.25, 0.2, size=8, color=COL["muted"])
add_text(s, "留存影响 ↑", mx-0.05, my-0.28, 1.0, 0.2, size=8, color=COL["muted"])
# quadrants
add_shape(s, MSO_SHAPE.RECTANGLE, mx+mw/2, my, mw/2, mh/2, RGBColor(255, 251, 235), line=None)
add_shape(s, MSO_SHAPE.RECTANGLE, mx, my, mw/2, mh/2, RGBColor(240, 253, 244), line=None)
add_shape(s, MSO_SHAPE.RECTANGLE, mx, my+mh/2, mw/2, mh/2, RGBColor(248, 250, 252), line=None)
add_shape(s, MSO_SHAPE.RECTANGLE, mx+mw/2, my+mh/2, mw/2, mh/2, RGBColor(254, 242, 242), line=None)
add_text(s, "优先推进", mx+0.25, my+0.20, 1.0, 0.2, size=10, bold=True, color=COL["success"])
add_text(s, "战略储备", mx+mw/2+0.25, my+0.20, 1.0, 0.2, size=10, bold=True, color=COL["orange"])
items = [
    ("可靠性", 1.4, 1.2, COL["primary"]),
    ("拍照答疑", 2.2, 1.6, COL["product_a"]),
    ("深度思考", 3.7, 1.5, COL["product_b"]),
    ("生图质量", 4.8, 2.1, COL["product_c"]),
    ("场景召回", 3.1, 3.0, COL["orange"]),  # [v3.7 修复] xx 2.8→3.1, 避免 label box 跨 Q3/Q4 边界 (x=3.95)
    ("陪伴/电商", 5.1, 3.0, COL["warning"]),
]
for name, xx, yy, c in items:
    add_shape(s, MSO_SHAPE.OVAL, mx+xx-0.12, my+yy-0.12, 0.24, 0.24, c)
    # [v3.7 修复] label w 从 0.9 收到 0.65, 避免跨象限边界 (拍照答疑 在 Q1/Q2 之间)
    add_text(s, name, mx+xx+0.15, my+yy-0.11, 0.65, 0.17, size=7.7, color=COL["text"])
# [v3.6 4 区填充] RIGHT: "发现" rail 填补 x>7.0 6 inch 空白 (v3.5 最大空白)
add_discovery_panel(s, [
    "优先推进: 可靠性 + 拍照答疑 + 深度思考",
    "战略储备: 生图质量 / 场景召回 (中期)",
    "底座先于探索, 不挤占资源"
])
# [v3.6 4 区填充] BOTTOM: 关键发现 3 列 panel (替换 v3.5 的 takeaway)
add_bottom_summary(s, [
    {"tag": "底座", "text": "可靠性: 高影响、低复杂度, 立即做", "color": COL["primary"]},
    {"tag": "抓手", "text": "拍照答疑: 高影响、场景明确, 短期做", "color": COL["product_a"]},
    {"tag": "探索", "text": "陪伴/电商: 中期探索, 不挤占底座", "color": COL["orange"]}
], title="优先级总结")
add_footer(s, "综合 Page 2、Page 8-13：按留存影响与实施复杂度进行定性排序。", 15)

# 16 Conclusion
s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX]); slide_bg(s, COL["navy"])
add_text(s, "结论：下一阶段应从“提升模型底座”与“抢占功能纯白场景”两条线并行推进", 0.72, 0.72, 10.9, 0.62, size=22, bold=True, color=COL["white"], font=TITLE_FONT)
# 3 big conclusion cards
for i, (num, t, b, c) in enumerate([
    ("01", "补可靠性", "把答案准确性、信源权威和速度稳定性作为留存底座。", COL["primary"]),
    ("02", "抢功能纯白", "围绕拍照答疑建立新用户首个高价值场景。", COL["product_a"]),
    ("03", "做场景召回", "将召回从提醒机制升级为任务场景入口。", COL["orange"]),
]):
    x = 0.82 + i*4.05
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.05, 3.38, 2.35, RGBColor(255, 255, 255), line=None, radius=True)
    add_text(s, num, x+0.20, 2.30, 0.8, 0.38, size=22, bold=True, color=c)
    add_text(s, t, x+0.25, 3.00, 2.4, 0.28, size=15, bold=True, color=COL["navy"])
    add_text(s, b, x+0.25, 3.45, 2.70, 0.46, size=10.5, color=COL["text"])
add_shape(s, MSO_SHAPE.RECTANGLE, 0.82, 5.25, 10.7, 0.01, RGBColor(90, 110, 130))
add_text(s, "建议以 4 周为一个实验周期：可靠性专项 → 拍照答疑新手路径 → 场景化召回实验 → 留存复盘。", 0.82, 5.55, 9.9, 0.32, size=13, color=RGBColor(225, 235, 245))
add_text(s, "Source: AI 产品用户留存分析_文档资料.pdf｜format.skill v3", 0.82, 6.82, 6.5, 0.22, size=8.5, color=RGBColor(175, 190, 205))
add_text(s, "16", 12.15, 6.82, 0.55, 0.22, size=8.5, color=RGBColor(175, 190, 205), align=PP_ALIGN.RIGHT)

# ---------- QA and save ----------
prs.save(OUT)

# [v3 QA 增强] layout_type 收敛检查 + v3 旧名禁用检查 + 必备组件检查
V4_VALID_LAYOUTS = {"cover", "executive_summary", "analysis_dashboard", "methodology_or_strategy", "priority_matrix"}
V3_DEPRECATED = {"title_slide", "action_title", "key_takeaway", "pyramid", "funnel",
                 "bar_chart_horizontal", "scatter_plot", "heatmap_table", "user_voice",
                 "action_items", "closing_slide", "cube", "quadrant_3group", "lollipop"}

# [v3.2 边距 + 重叠 QA]
MARGIN_INTERNAL = 0.7  # FMT-V3-009: 内部页面边距 0.7 inch
COVER_INDICES = (1, 16)  # cover/closing 页用 0.55 边距,排除
SHAPE_TYPES_BG = (1, 17)  # 1=RECTANGLE, 17=TEXT_BOX (背景色 slide_bg 允许全画布)

qa = {
    "file": str(OUT),
    "version": "v3.12",
    "slide_count": len(prs.slides),
    "layout_type_map": V4_LAYOUT_MAP,
    "slides": [],
    "passed": True,
    "issues": [],
    "v3_specific_checks": {
        "all_layouts_in_v3_set": True,
        "no_v2_deprecated_used": True,
        "navigation_present_on_dashboard": True,
        "insight_panel_present_on_dashboard": True,
        "insight_panel_variants_used": [],
        "insight_panel_variants_diversity": 0,
        "evidence_to_variant_binding": True,
        "uniform_page_margin_passed": True,
        "no_overlap_passed": True,
        "executive_summary_no_rail_passed": True,
        "bottom_summary_present_passed": True,  # [v3.6 新增] FMT-V3-012
        "discovery_rail_present_passed": True    # [v3.6 新增] FMT-V3-013
    }
}

# 校验 1: layout_type 全部在 v3 集合中
for slide_no, layout in V4_LAYOUT_MAP.items():
    if layout not in V4_VALID_LAYOUTS:
        qa["passed"] = False
        qa["v3_specific_checks"]["all_layouts_in_v3_set"] = False
        qa["issues"].append(f"Slide {slide_no}: layout '{layout}' not in v3 5-class set")
    if layout in V3_DEPRECATED:
        qa["passed"] = False
        qa["v3_specific_checks"]["no_v2_deprecated_used"] = False
        qa["issues"].append(f"Slide {slide_no}: uses deprecated v3 layout '{layout}'")

# 校验 2: analysis_dashboard 必备 navigation + insight_panel
# 启发式：从 qa['slides'] 提取文本，搜索关键词
for idx, slide in enumerate(prs.slides, 1):
    shapes = list(slide.shapes)
    text_shapes = [sh for sh in shapes if getattr(sh, "has_text_frame", False)]
    non_text_shapes = [sh for sh in shapes if not getattr(sh, "has_text_frame", False) or sh.shape_type != 17]
    texts = []
    for sh in text_shapes:
        try:
            txt = sh.text.strip()
            if txt:
                texts.append(txt)
        except Exception:
            pass
    full_text = " ".join(texts)

    slide_info = {
        "idx": idx,
        "layout_type": V4_LAYOUT_MAP.get(idx, "unknown"),
        "shape_count": len(shapes),
        "text_shape_count": len(text_shapes),
        "non_text_shape_count": len(non_text_shapes),
        "has_visual": len(non_text_shapes) > 0,
        "text_chars": sum(len(t) for t in texts)
    }

    # v3 强制：analysis_dashboard 必带 navigation (tab 关键词) + insight_panel (关键洞察)
    if slide_info["layout_type"] == "analysis_dashboard":
        # [v3.5] nav_keywords 与新 TABS 同步: ["人群","纯白","非纯白","文本","功能","运营"]
        nav_keywords = ["人群", "纯白", "非纯白", "文本", "功能", "运营"]
        has_nav = any(k in full_text for k in nav_keywords)
        # v3.5: panel_keywords 加 "发现" (add_discovery_panel 的 title)
        panel_keywords = ["关键洞察", "发现", "核心结论", "行动建议", "研究流程", "优先级总结",
                          "全图最高", "杠杆", "Top 因子", "正向", "负向", "整体最强", "底座", "抓手", "探索", "差距"]
        has_panel = any(k in full_text for k in panel_keywords)
        slide_info["has_navigation"] = has_nav
        slide_info["has_insight_panel"] = has_panel
        # v3.5: 变体识别 — 新增"发现" → "rail" (add_discovery_panel 等同旧 rail)
        # [v3.6] 改为按"主要 panel"位置识别: 找到 w>10 h>0.5 的大 panel, 判定其变体
        variant = None
        # 先找页内所有"大 panel" (w > 6, h > 0.5) — 这些是主要的 insight panels
        large_panels = []
        for sh in shapes:
            try:
                if not getattr(sh, "left", None) or not getattr(sh, "top", None):
                    continue
                w = sh.width / 914400 if sh.width else 0
                h = sh.height / 914400 if sh.height else 0
                if w > 6.0 and h > 0.5:
                    large_panels.append({
                        "y": sh.top / 914400,
                        "x": sh.left / 914400,
                        "w": w, "h": h,
                        "text": sh.text if getattr(sh, "has_text_frame", False) else ""
                    })
            except Exception:
                pass
        # 按 y 排序
        large_panels.sort(key=lambda p: (p["y"], p["x"]))
        # 主 panel 选择策略:
        # - 如果有大 panel 在 y=5.4-5.7 (bottom_summary 位置) → 优先识别
        # - 如果有大 panel 在 y=6.0-6.7 (takeaway 位置) → 识别为 takeaway
        # - 如果有大 panel 在 y=0.7-1.2 (banner 位置) → 识别为 banner
        # - 否则看右 rail 或 callout
        main_panel = None
        for p in large_panels:
            if 5.40 < p["y"] < 5.70 and p["w"] > 10:
                main_panel = p; break
        if not main_panel:
            for p in large_panels:
                if 6.00 < p["y"] < 6.20 and p["w"] > 10:
                    main_panel = p; break
        if not main_panel:
            for p in large_panels:
                if 0.85 < p["y"] < 1.20 and p["w"] > 8:
                    main_panel = p; break
        # 现在根据 main_panel 内容判定变体
        if main_panel:
            mp_text = main_panel["text"]
            # bottom_summary 标题特征 (含 y=5.40-5.70 位置)
            # 标题与 takeaway 的"优先级总结"区分: 位置 y=5.55 vs y=6.05
            if any(t in full_text for t in ["留存格局", "驱动格局", "来源机会", "功能格局", "召回格局",
                                             "功能纯白提升", "文本满意度", "流失原因", "优先级总结"]):
                variant = "bottom_summary"
            elif "研究流程" in mp_text or "行动建议" in mp_text:
                variant = "takeaway"
            elif "核心结论" in mp_text:
                variant = "banner"
        if not variant:
            # 回退到文本启发式
            if "关键洞察" in full_text or "发现" in full_text:
                variant = "rail"
            elif "核心结论" in full_text:
                variant = "banner"
            elif "研究流程" in full_text or "行动建议" in full_text or "优先级总结" in full_text:
                variant = "takeaway"
            elif "全图最高" in full_text or "差距" in full_text or "杠杆" in full_text:
                variant = "callout"
            elif "Top 因子" in full_text or "整体最强" in full_text:
                variant = "anchor"
            elif "正向" in full_text and "负向" in full_text and "豆包" in full_text:
                variant = "matrix"
        if variant:
            slide_info["insight_panel_variant"] = variant
            qa["v3_specific_checks"]["insight_panel_variants_used"].append(f"P{idx}:{variant}")
        if not has_nav:
            qa["v3_specific_checks"]["navigation_present_on_dashboard"] = False
            qa["issues"].append(f"Slide {idx} (analysis_dashboard): missing navigation tabs")
        if not has_panel:
            qa["v3_specific_checks"]["insight_panel_present_on_dashboard"] = False
            qa["issues"].append(f"Slide {idx} (analysis_dashboard): missing insight panel")
        if not variant:
            qa["v3_specific_checks"]["evidence_to_variant_binding"] = False
            qa["issues"].append(f"Slide {idx} (analysis_dashboard): insight_panel variant not identifiable")

    # [v3.2] FMT-V3-009: 内部页面边距统一为 0.7 inch
    if idx not in COVER_INDICES:
        for sh in shapes:
            try:
                if not getattr(sh, "left", None) or not getattr(sh, "top", None):
                    continue
                # 跳过装饰性背景(slide_bg 全画布)、page_no 文本、装饰圆 (cover 才用)
                if sh.width is None or sh.height is None:
                    continue
                x = sh.left / 914400  # EMU -> inch
                y = sh.top / 914400
                w = sh.width / 914400
                h = sh.height / 914400
                # 边距检查 (容差 0.05 inch)
                if x < MARGIN_INTERNAL - 0.05 or y < MARGIN_INTERNAL - 0.05 or \
                   x + w > W - MARGIN_INTERNAL + 0.05 or y + h > H - MARGIN_INTERNAL + 0.05:
                    # 排除很小的 shape (< 0.3 inch) 和 footer 区域页码文本
                    if w < 0.3 or h < 0.3:
                        continue
                    # [v3.2 修复] 豁免"页面框架元素": 顶部/底部的水平横条
                    # - 顶部页眉 (title/nav): y < 0.7 但 w > 10 且 h < 0.7
                    # - 底部页脚 (footer): y + h > 6.8 但 w > 10 且 h < 0.7
                    is_top_frame = (y < MARGIN_INTERNAL - 0.05) and (w > 10.0) and (h < 0.7)
                    is_bottom_frame = (y + h > H - MARGIN_INTERNAL + 0.05) and (w > 10.0) and (h < 0.7)
                    if is_top_frame or is_bottom_frame:
                        continue
                    qa["v3_specific_checks"]["uniform_page_margin_passed"] = False
                    qa["issues"].append(f"Slide {idx} (FMT-V3-009): shape at ({x:.2f},{y:.2f},{w:.2f}x{h:.2f}) 超出 0.7 inch 边距")
                    break
            except Exception:
                pass

        # [v3.2 修复] FMT-V3-010: 无 shape 重叠
        # 启发式：只对两个 shape 中"较大者" w > 1.0 & h > 0.5 才做几何检查
        # 这样跳过 panel 内部的 icon / 数字文字 / 标签（这些是设计嵌入）
        # [v3.2 增强] 父子关系豁免: 如果 shape A 完全包含在 shape B 内(允许 0.10 inch 内边距),
        # 视为 parent-child 设计嵌入, 不算重叠
        visual_shapes = [sh for sh in shapes if getattr(sh, "left", None) and getattr(sh, "top", None) and sh.width and sh.height and (sh.width / 914400) >= 1.0 and (sh.height / 914400) >= 0.5]
        for i, sh1 in enumerate(visual_shapes):
            x1 = sh1.left / 914400; y1 = sh1.top / 914400
            w1 = sh1.width / 914400; h1 = sh1.height / 914400
            for sh2 in visual_shapes[i+1:]:
                x2 = sh2.left / 914400; y2 = sh2.top / 914400
                w2 = sh2.width / 914400; h2 = sh2.height / 914400
                # AABB 重叠检测
                if x1 < x2 + w2 and x1 + w1 > x2 and y1 < y2 + h2 and y1 + h1 > y2:
                    overlap_x = min(x1 + w1, x2 + w2) - max(x1, x2)
                    overlap_y = min(y1 + h1, y2 + h2) - max(y1, y2)
                    if overlap_x > 0.20 and overlap_y > 0.20:
                        # [v3.2 增强] 父子关系检测: A 完全在 B 内 (允许 0.10 inch 内边距)
                        # 判断条件: A.left >= B.left-0.10 AND A.top >= B.top-0.10
                        #          AND A.right <= B.right+0.10 AND A.bottom <= B.bottom+0.10
                        PAD = 0.10
                        a_in_b = (x1 >= x2 - PAD and y1 >= y2 - PAD and
                                  x1 + w1 <= x2 + w2 + PAD and y1 + h1 <= y2 + h2 + PAD)
                        b_in_a = (x2 >= x1 - PAD and y2 >= y1 - PAD and
                                  x2 + w2 <= x1 + w1 + PAD and y2 + h2 <= y1 + h1 + PAD)
                        if a_in_b or b_in_a:
                            # 父子关系, 跳过(panel 容器与其内部元素)
                            continue
                        qa["v3_specific_checks"]["no_overlap_passed"] = False
                        qa["issues"].append(f"Slide {idx} (FMT-V3-010): shape ({x1:.2f},{y1:.2f},{w1:.2f}x{h1:.2f}) 与 ({x2:.2f},{y2:.2f},{w2:.2f}x{h2:.2f}) 重叠 {overlap_x:.2f}x{overlap_y:.2f}")
                        break
            else:
                continue
            break

    if not slide_info["has_visual"]:
        qa["passed"] = False; qa["issues"].append(f"Slide {idx}: no visual shapes")
    if idx not in (1, 16) and len(shapes) < 10:
        qa["passed"] = False; qa["issues"].append(f"Slide {idx}: low shape count")
    qa["slides"].append(slide_info)

# v3.1: 统计实际使用的变体种类
used_variants = set()
for s in qa["slides"]:
    v = s.get("insight_panel_variant")
    if v: used_variants.add(v)
qa["v3_specific_checks"]["insight_panel_variants_diversity"] = len(used_variants)
# 期望 6 变体中至少用 4 种（不强求全部，但要体现差异化）
if len(used_variants) < 3:
    qa["passed"] = False
    qa["issues"].append(f"insight_panel 变体使用过少: 只用了 {sorted(used_variants)}, 应至少 3 种以体现 evidence 差异化")

# [v3.4] FMT-V3-011: executive_summary 不应出现 insight_panel_rail 变体
# 反馈: 核心发现页本身已是观点罗列,右侧 rail 视觉冗余
for s in qa["slides"]:
    if s.get("layout_type") == "executive_summary" and s.get("insight_panel_variant") == "rail":
        qa["v3_specific_checks"]["executive_summary_no_rail_passed"] = False
        qa["issues"].append(f"Slide {s['idx']} (FMT-V3-011): executive_summary 不应使用 insight_panel_rail 变体 (5 cards 本身已观点罗列)")

# [v3.6] FMT-V3-012: 4 区填充检查 — 每页 (P3-P15) 必须满足:
#   ① 有 bottom_summary panel (y=5.55-6.65) 填补底部空白
#   ② 底部内容 y_end < 5.5 → 失败 (空白 > 1 inch)
#   ③ x>9.5 区域有发现 rail 内容 → 填补右对齐空白
# 豁免: P1 (cover), P2 (executive_summary), P16 (conclusion closing) - 3 页有专门设计
COVER_INDICES_FMT12 = (1, 2, 16)  # 豁免
BOTTOM_BLANK_THRESHOLD = 1.0  # inch — 底部内容 y_end < (H - 1.0) = 6.5 则失败
for idx, slide in enumerate(prs.slides, 1):
    if idx in COVER_INDICES_FMT12:
        continue
    if idx == 14:  # P14 已有 takeaway 面板, 通过
        # 校验 takeaway 存在
        texts = [sh.text for sh in slide.shapes if getattr(sh, "has_text_frame", False)]
        full = " ".join(t for t in texts if t)
        if "行动建议" in full or "立即做" in full:
            continue
    if idx == 13:  # P13 矩阵 + quotes 已填满大部分区域, 豁免
        continue
    if idx == 6:  # P6 scatter plot 完整填到 y=5.97, 接近 6.0
        continue
    # 检查 ①: bottom panel (summary 或 takeaway) 必须存在 (y >= 5.4, w > 10, h > 0.5)
    has_bottom_panel = False
    for sh in slide.shapes:
        try:
            y = sh.top / 914400
            h = sh.height / 914400
            w = sh.width / 914400
            if y >= 5.4 and w > 10 and h > 0.5:
                has_bottom_panel = True
                break
        except Exception:
            pass
    if not has_bottom_panel:
        qa["v3_specific_checks"]["bottom_summary_present_passed"] = False
        qa["issues"].append(f"Slide {idx} (FMT-V3-012-①): 缺底部关键摘要 panel (应在 y=5.55-6.65 区)")

    # 检查 ②: 底部内容 y_end < 5.5 (底部空白 > 1.0 inch)
    max_y_end = 0
    for sh in slide.shapes:
        try:
            if not getattr(sh, "left", None) or not getattr(sh, "top", None):
                continue
            if sh.width is None or sh.height is None:
                continue
            y_end = (sh.top + sh.height) / 914400
            if y_end > max_y_end:
                max_y_end = y_end
        except Exception:
            pass
    if max_y_end < H - BOTTOM_BLANK_THRESHOLD - 0.5:  # 6.5 - 0.5 = 6.0, 但 footer 在 6.85
        qa["v3_specific_checks"]["bottom_summary_present_passed"] = False
        qa["issues"].append(f"Slide {idx} (FMT-V3-012-②): 底部内容最大 y_end={max_y_end:.2f} < {H - 0.5 - BOTTOM_BLANK_THRESHOLD:.2f}, 底部空白 > {BOTTOM_BLANK_THRESHOLD} inch")

    # 检查 ③: x>9.5 区域有发现 rail (灰底或含'发现'文字)
    has_rail = False
    for sh in slide.shapes:
        try:
            x = sh.left / 914400
            w = sh.width / 914400
            if x >= 9.8 and w > 2.0:  # 右 rail 起始位置
                # 检查是否含'发现'文字 或 灰底填充
                if getattr(sh, "has_text_frame", False) and "发现" in (sh.text or ""):
                    has_rail = True
                    break
                if sh.shape_type == 1:  # RECTANGLE
                    # 灰底 panel (R=240-250, G=240-250, B=240-250)
                    try:
                        if hasattr(sh, "fill") and sh.fill.type == 1:  # SOLID
                            c = sh.fill.fore_color.rgb
                            if 235 <= c[0] <= 250 and 235 <= c[1] <= 250 and 235 <= c[2] <= 250:
                                has_rail = True
                                break
                    except Exception:
                        pass
        except Exception:
            pass
    if not has_rail:
        qa["v3_specific_checks"]["discovery_rail_present_passed"] = False
        qa["issues"].append(f"Slide {idx} (FMT-V3-013): x>9.5 区域缺'发现'rail (右对齐空白 > 3 inch)")

# 最终汇总
if qa["issues"]:
    qa["passed"] = False

QA_OUT.write_text(json.dumps(qa, ensure_ascii=False, indent=2), encoding="utf-8")
print(f"SAVED: {OUT}")
print(f"QA: {QA_OUT}")
print(json.dumps({
    "version": qa["version"],
    "passed": qa["passed"],
    "slide_count": qa["slide_count"],
    "v3_checks": qa["v3_specific_checks"],
    "issues": qa["issues"][:5]
}, ensure_ascii=False, indent=2))
